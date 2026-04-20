from __future__ import annotations
import os, sys, httpx, sqlite3, json, logging, re, base64, io, random, asyncio, time
from logging.handlers import RotatingFileHandler
from collections import defaultdict
from datetime import datetime, timedelta, timezone
from telegram import Update, ReplyKeyboardMarkup, KeyboardButton, InlineKeyboardMarkup, InlineKeyboardButton
from telegram.ext import Application, MessageHandler, CommandHandler, CallbackQueryHandler, filters, ContextTypes
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build
from google.oauth2.credentials import Credentials

# Большие константы-списки (цитаты дня, вопросы рефлексии) вынесены
# в отдельный файл prompts.py — раньше они занимали 160+ строк в bot.py.
from prompts import QUOTES, REFLECT_QUESTIONS
# UI-клавиатуры и константа SPHERES. В bot.py остались только клавиатуры,
# которые тянут данные из БД (habits_keyboard, settings_keyboard).
from keyboards import (
    SPHERES, SPHERE_KEYS,
    score_keyboard, main_keyboard, onboarding_keyboard,
    tasks_keyboard, goals_keyboard, spheres_keyboard, sphere_detail_keyboard,
    task_actions_keyboard, move_timeframe_keyboard,
)

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

# Логирование: консоль + файл с ротацией (10 МБ × 3 файла = 30 МБ лимит).
# Ошибки всегда видно в nova.log — можно читать прямо на Railway.
_log_fmt = "%(asctime)s %(levelname)s [%(name)s] %(message)s"
_root_logger = logging.getLogger()
_root_logger.setLevel(logging.INFO)
if not _root_logger.handlers:
    _stream = logging.StreamHandler(sys.stdout)
    _stream.setFormatter(logging.Formatter(_log_fmt))
    _root_logger.addHandler(_stream)
    try:
        _file = RotatingFileHandler("nova.log", maxBytes=10_000_000, backupCount=3, encoding="utf-8")
        _file.setFormatter(logging.Formatter(_log_fmt))
        _root_logger.addHandler(_file)
    except Exception as _e:
        logging.warning(f"File logging disabled: {_e}")
# Снижаем шум сторонних библиотек
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)
logging.getLogger("telegram").setLevel(logging.WARNING)
logging.getLogger("apscheduler").setLevel(logging.WARNING)

# Вся конфигурация — секреты, тарифы, модели, лимиты токенов, ключевые слова —
# вынесена в config.py. Правь цены и лимиты там.
from config import (
    TELEGRAM_TOKEN, CLAUDE_API_KEY, TURSO_URL, TURSO_TOKEN,
    GOOGLE_CLIENT_ID, GOOGLE_CLIENT_SECRET, WEBHOOK_URL,
    OPENROUTER_API_KEY, MEM0_API_KEY, OWNER_ID, REPLICATE_API_TOKEN,
    SCOPES, PAYMENTS_ENABLED, PLANS, PLAN_EXPIRED,
    MODEL_SMART, MODEL_FAST_OPENROUTER, MODEL_FAST_CLAUDE,
    SMART_KEYWORDS as _SMART_KEYWORDS,
    MAX_TOKENS_DEFAULT, MAX_TOKENS_ONBOARD, MAX_TOKENS_NOTIF,
    MAX_TOKENS_REVIEW, MAX_TOKENS_VISION,
)
# Обёртки над внешними сервисами — Anthropic/OpenRouter/Groq/Mem0.
from services.claude import pick_model, call_claude, call_claude_vision
from services.mem0 import get_mem0, mem0_add, mem0_search, mem0_delete_all_user
from services.voice import call_groq_voice


_SQLITE_PRAGMA_APPLIED = False

def _apply_sqlite_pragma(conn):
    """WAL-режим: параллельные чтения не блокируют записи, меньше I/O.
    Применяем один раз за процесс — это настройки уровня файла БД."""
    global _SQLITE_PRAGMA_APPLIED
    if _SQLITE_PRAGMA_APPLIED:
        return
    try:
        c = conn.cursor()
        c.execute("PRAGMA journal_mode=WAL")
        c.execute("PRAGMA synchronous=NORMAL")
        c.execute("PRAGMA temp_store=MEMORY")
        c.execute("PRAGMA cache_size=-20000")  # 20 МБ кэша страниц
        c.execute("PRAGMA busy_timeout=5000")  # ждём до 5с если БД занята
        conn.commit()
        _SQLITE_PRAGMA_APPLIED = True
    except Exception as e:
        logging.warning(f"SQLite PRAGMA failed: {e}")

def get_conn(sync=True):
    if TURSO_URL and TURSO_TOKEN:
        try:
            import libsql_experimental as libsql
            conn = libsql.connect("nova.db", sync_url=TURSO_URL, auth_token=TURSO_TOKEN)
            if sync:
                conn.sync()
            return conn
        except Exception as e:
            logging.warning(f"Turso failed: {e}")
    conn = sqlite3.connect("assistant.db", timeout=5)
    _apply_sqlite_pragma(conn)
    return conn

def init_db():
    conn = get_conn()
    c = conn.cursor()
    c.execute("""CREATE TABLE IF NOT EXISTS users (
        user_id INTEGER PRIMARY KEY,
        onboarding_done INTEGER DEFAULT 0,
        onboarding_step INTEGER DEFAULT 0,
        profile TEXT DEFAULT '{}')""")
    c.execute("""CREATE TABLE IF NOT EXISTS messages (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, role TEXT, content TEXT,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP)""")
    c.execute("""CREATE TABLE IF NOT EXISTS tasks (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, text TEXT,
        priority TEXT DEFAULT 'normal',
        sphere TEXT DEFAULT 'general',
        timeframe TEXT DEFAULT 'week',
        done INTEGER DEFAULT 0,
        due_date TEXT, created_at TEXT,
        done_at TEXT)""")
    try:
        c.execute("ALTER TABLE tasks ADD COLUMN done_at TEXT")
    except Exception:
        pass
    c.execute("""CREATE TABLE IF NOT EXISTS goals (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, text TEXT,
        sphere TEXT DEFAULT 'general',
        timeframe TEXT DEFAULT 'longterm',
        progress INTEGER DEFAULT 0,
        done INTEGER DEFAULT 0, created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS ideas (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, text TEXT,
        sphere TEXT DEFAULT 'general',
        reviewed_at TEXT,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS sphere_activity (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER, sphere TEXT, activity_date TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS google_tokens (
        user_id INTEGER PRIMARY KEY,
        token TEXT,
        refresh_token TEXT,
        token_uri TEXT,
        client_id TEXT,
        client_secret TEXT,
        scopes TEXT,
        expiry TEXT)""")
    try:
        c.execute("ALTER TABLE google_tokens ADD COLUMN expiry TEXT")
    except Exception:
        pass
    c.execute("""CREATE TABLE IF NOT EXISTS sent_quotes (
        user_id INTEGER,
        quote_idx INTEGER,
        PRIMARY KEY (user_id, quote_idx))""")
    c.execute("""CREATE TABLE IF NOT EXISTS followup_queue (
        user_id INTEGER PRIMARY KEY,
        asked_at TEXT,
        attempts INTEGER DEFAULT 0)""")
    c.execute("""CREATE TABLE IF NOT EXISTS mood_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        score INTEGER,
        note TEXT,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS energy_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        score INTEGER,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS habits (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        name TEXT,
        active INTEGER DEFAULT 1,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS habit_log (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        habit_id INTEGER,
        log_date TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS journal (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        question TEXT,
        entry TEXT,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS wins (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        text TEXT,
        created_at TEXT)""")
    # Подписка юзера: какой тариф, до какой даты действует.
    c.execute("""CREATE TABLE IF NOT EXISTS subscriptions (
        user_id INTEGER PRIMARY KEY,
        plan TEXT DEFAULT 'free',
        valid_until TEXT,
        last_payment_stars INTEGER,
        last_payment_at TEXT)""")
    # Счётчик использования по дням. Ключ (user_id, day) — один ряд на юзера в сутки.
    c.execute("""CREATE TABLE IF NOT EXISTS usage_counters (
        user_id INTEGER,
        day TEXT,
        kind TEXT,
        count INTEGER DEFAULT 0,
        PRIMARY KEY (user_id, day, kind))""")
    # Распознанные чеки / траты (из фото или ручной ввод).
    c.execute("""CREATE TABLE IF NOT EXISTS expenses (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        amount REAL,
        currency TEXT DEFAULT 'RUB',
        category TEXT,
        note TEXT,
        created_at TEXT)""")
    # «Парковочная площадка» — темы, которые человек затронул в онбординге,
    # но для которых сейчас не время. После завершения онбординга Нова
    # вернётся к ним по запросу пользователя (кнопка «Узнай меня больше»
    # или команда /parking).
    c.execute("""CREATE TABLE IF NOT EXISTS parking_lot (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        topic TEXT,
        created_at TEXT,
        discussed_at TEXT)""")
    # Рефералы: кто кого пригласил. Заполняется при /start <ref_...>
    c.execute("""CREATE TABLE IF NOT EXISTS referrals (
        invited_user_id INTEGER PRIMARY KEY,
        inviter_user_id INTEGER,
        created_at TEXT,
        rewarded_at TEXT)""")
    # Промокоды: разовая активация даёт бонусные дни подписки.
    c.execute("""CREATE TABLE IF NOT EXISTS promo_codes (
        code TEXT PRIMARY KEY,
        plan TEXT,
        days INTEGER,
        max_uses INTEGER DEFAULT 1,
        uses INTEGER DEFAULT 0,
        expires_at TEXT,
        created_at TEXT)""")
    c.execute("""CREATE TABLE IF NOT EXISTS promo_redemptions (
        user_id INTEGER,
        code TEXT,
        redeemed_at TEXT,
        PRIMARY KEY (user_id, code))""")
    # NPS/отзыв: через 14 дней после старта Нова спрашивает «Как тебе?».
    # Ответ сохраняется. Поле asked_at = когда спросили, rating/feedback = ответ.
    c.execute("""CREATE TABLE IF NOT EXISTS feedback (
        user_id INTEGER PRIMARY KEY,
        asked_at TEXT,
        rating INTEGER,
        text TEXT,
        answered_at TEXT)""")

    # Индексы для всех горячих выборок по user_id — без них
    # SQLite сканирует всю таблицу, что на 500+ юзерах даёт задержки в секундах.
    for idx_sql in [
        "CREATE INDEX IF NOT EXISTS idx_messages_user      ON messages(user_id, id DESC)",
        "CREATE INDEX IF NOT EXISTS idx_tasks_user         ON tasks(user_id, done)",
        "CREATE INDEX IF NOT EXISTS idx_tasks_user_due     ON tasks(user_id, due_date)",
        "CREATE INDEX IF NOT EXISTS idx_goals_user         ON goals(user_id, done)",
        "CREATE INDEX IF NOT EXISTS idx_ideas_user         ON ideas(user_id)",
        "CREATE INDEX IF NOT EXISTS idx_sphere_user_date   ON sphere_activity(user_id, activity_date)",
        "CREATE INDEX IF NOT EXISTS idx_mood_user          ON mood_log(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_energy_user        ON energy_log(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_habits_user        ON habits(user_id, active)",
        "CREATE INDEX IF NOT EXISTS idx_habitlog_user_date ON habit_log(user_id, log_date)",
        "CREATE INDEX IF NOT EXISTS idx_journal_user       ON journal(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_wins_user          ON wins(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_usage_user_day     ON usage_counters(user_id, day)",
        "CREATE INDEX IF NOT EXISTS idx_expenses_user      ON expenses(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_parking_user       ON parking_lot(user_id, discussed_at)",
        "CREATE INDEX IF NOT EXISTS idx_referrals_inviter  ON referrals(inviter_user_id)",
    ]:
        try:
            c.execute(idx_sql)
        except Exception as e:
            logging.warning(f"Index skip: {e}")

    conn.commit()
    if hasattr(conn, 'sync'): conn.sync()
    conn.close()

# Таблицы, которые принадлежат конкретному юзеру.
# Используется в /reset, /newuser и "полном сбросе" из диалога —
# раньше этот список дублировался в трёх местах.
USER_DATA_TABLES = [
    "users", "messages", "tasks", "goals", "ideas",
    "sphere_activity", "google_tokens",
    "mood_log", "energy_log", "habits", "habit_log",
    "journal", "wins", "sent_quotes", "followup_queue",
    "subscriptions", "usage_counters", "expenses", "parking_lot",
    "referrals", "promo_redemptions", "feedback",
]

def wipe_user_data(uid):
    """Полное удаление всех данных пользователя. Используется в /reset, /newuser
    и ручном сбросе через диалог."""
    for t in USER_DATA_TABLES:
        try:
            db_exec(f"DELETE FROM {t} WHERE user_id=?", (uid,))
        except Exception as e:
            logging.warning(f"wipe {t} for {uid}: {e}")

def db_exec(query, params=()):
    conn = get_conn()
    c = conn.cursor()
    c.execute(query, tuple(params))
    conn.commit()
    if hasattr(conn, 'sync'): conn.sync()
    conn.close()

def db_fetch(query, params=()):
    conn = get_conn(sync=False)
    c = conn.cursor()
    c.execute(query, tuple(params))
    rows = c.fetchall()
    conn.close()
    return rows

def db_fetchone(query, params=()):
    conn = get_conn(sync=False)
    c = conn.cursor()
    c.execute(query, tuple(params))
    row = c.fetchone()
    conn.close()
    return row

def ensure_user(uid):
    db_exec("INSERT OR IGNORE INTO users (user_id) VALUES (?)", (uid,))
    # Одновременно заводим триальную подписку на 7 дней, если юзер новый.
    # Функция ensure_subscription определена ниже — это ок, вызов в runtime.
    try:
        ensure_subscription(uid)
    except NameError:
        # Защита от случая, когда ensure_user вызывается до определения ensure_subscription
        # (не должно случаться, но пусть будет).
        pass

def get_user(uid):
    return db_fetchone("SELECT * FROM users WHERE user_id=?", (uid,))

def update_user(uid, **kw):
    for k, v in kw.items():
        db_exec(f"UPDATE users SET {k}=? WHERE user_id=?", (v, uid))

def get_profile(uid):
    row = db_fetchone("SELECT profile FROM users WHERE user_id=?", (uid,))
    return json.loads(row[0]) if row else {}

def save_profile(uid, profile):
    db_exec("UPDATE users SET profile=? WHERE user_id=?",
            (json.dumps(profile, ensure_ascii=False), uid))

def save_msg(uid, role, content):
    db_exec("INSERT INTO messages (user_id,role,content,created_at) VALUES (?,?,?,?)",
            (uid, role, content, datetime.now().isoformat()))

def get_history(uid, limit=20):
    rows = db_fetch("SELECT role,content FROM messages WHERE user_id=? ORDER BY id DESC LIMIT ?", (uid, limit))
    return [{"role": r[0], "content": r[1]} for r in reversed(rows)]

def clear_history(uid):
    db_exec("DELETE FROM messages WHERE user_id=?", (uid,))

def add_task(uid, text, priority="normal", sphere="general", timeframe="week", due_date=None):
    existing = db_fetch("SELECT id FROM tasks WHERE user_id=? AND text=? AND done=0", (uid, text))
    if existing: return
    db_exec("INSERT INTO tasks (user_id,text,priority,sphere,timeframe,due_date,created_at) VALUES (?,?,?,?,?,?,?)",
            (uid, text, priority, sphere, timeframe, due_date, datetime.now().isoformat()))

def get_tasks(uid, sphere=None, timeframe=None, priority=None, done=0):
    query = "SELECT id,text,priority,sphere,timeframe,due_date FROM tasks WHERE user_id=? AND done=?"
    params = [uid, done]
    if sphere: query += " AND sphere=?"; params.append(sphere)
    if timeframe: query += " AND timeframe=?"; params.append(timeframe)
    if priority: query += " AND priority=?"; params.append(priority)
    query += " ORDER BY id"
    return db_fetch(query, tuple(params))

def get_today_tasks(uid):
    today = datetime.now().date().isoformat()
    return db_fetch("""SELECT id,text,priority,sphere,timeframe,due_date FROM tasks
                 WHERE user_id=? AND done=0 AND (timeframe='today' OR due_date=?)
                 ORDER BY priority DESC""", (uid, today))

def complete_task(task_id):
    db_exec("UPDATE tasks SET done=1, done_at=? WHERE id=?",
            (datetime.now().isoformat(), task_id))

def delete_task(task_id):
    db_exec("DELETE FROM tasks WHERE id=?", (task_id,))

_UNSET = object()

def edit_task(task_id, text=None, priority=None, timeframe=None, due_date=_UNSET):
    if text: db_exec("UPDATE tasks SET text=? WHERE id=?", (text, task_id))
    if priority: db_exec("UPDATE tasks SET priority=? WHERE id=?", (priority, task_id))
    if timeframe: db_exec("UPDATE tasks SET timeframe=? WHERE id=?", (timeframe, task_id))
    if due_date is not _UNSET: db_exec("UPDATE tasks SET due_date=? WHERE id=?", (due_date, task_id))

def update_goal_progress(goal_id, progress):
    db_exec("UPDATE goals SET progress=? WHERE id=?", (progress, goal_id))

def get_user_tz_offset(profile):
    tz = profile.get("timezone", "")
    if not tz:
        return 0
    m = re.search(r'([+-]?\d+)', tz)
    if m:
        try:
            return int(m.group(1))
        except Exception:
            pass
    return 0

def user_now(profile):
    return datetime.now(timezone.utc).replace(tzinfo=None) + timedelta(hours=get_user_tz_offset(profile))

# ── Цитаты ────────────────────────────────────────────────────────────────────

def get_random_quote(uid):
    sent = {r[0] for r in db_fetch("SELECT quote_idx FROM sent_quotes WHERE user_id=?", (uid,))}
    available = [i for i in range(len(QUOTES)) if i not in sent]
    if not available:
        db_exec("DELETE FROM sent_quotes WHERE user_id=?", (uid,))
        available = list(range(len(QUOTES)))
    idx = random.choice(available)
    db_exec("INSERT OR IGNORE INTO sent_quotes (user_id, quote_idx) VALUES (?,?)", (uid, idx))
    return QUOTES[idx]

# ── Rate-limit ────────────────────────────────────────────────────────────────
# Защита от потока сообщений от одного юзера (случайный цикл у клиента, троллинг,
# баг в интеграции). Окно: не более 20 сообщений за 60 секунд + минимум 0.4с
# между сообщениями. Цифры подобраны так, чтобы нормальный человек никогда
# не упёрся в лимит, но массовый спам был отсечён до похода в LLM.

_RATE_WINDOW_SEC = 60
_RATE_MAX_IN_WINDOW = 20
_RATE_MIN_GAP_SEC = 0.4
_rate_buckets: dict[int, list[float]] = defaultdict(list)
_rate_last: dict[int, float] = {}

def check_rate_limit(uid: int) -> tuple[bool, str]:
    """Возвращает (allowed, reason). reason — причина отказа, если allowed=False."""
    now = time.time()
    last = _rate_last.get(uid, 0.0)
    if now - last < _RATE_MIN_GAP_SEC:
        return False, "too_fast"
    bucket = _rate_buckets[uid]
    cutoff = now - _RATE_WINDOW_SEC
    # Чистим старые записи
    while bucket and bucket[0] < cutoff:
        bucket.pop(0)
    if len(bucket) >= _RATE_MAX_IN_WINDOW:
        return False, "too_many"
    bucket.append(now)
    _rate_last[uid] = now
    return True, ""

# ── Follow-up ─────────────────────────────────────────────────────────────────

def set_followup(uid):
    db_exec("INSERT OR REPLACE INTO followup_queue (user_id, asked_at, attempts) VALUES (?,?,0)",
            (uid, datetime.now().isoformat()))

def clear_followup(uid):
    db_exec("DELETE FROM followup_queue WHERE user_id=?", (uid,))

def get_pending_followups():
    threshold = (datetime.now() - timedelta(hours=1)).isoformat()
    return db_fetch("SELECT user_id, asked_at, attempts FROM followup_queue WHERE asked_at < ? AND attempts < 2",
                    (threshold,))

# ── Parking lot (отложенные темы из онбординга) ──────────────────────────────

def add_parking(uid: int, topic: str):
    """Отметить, что человек затронул тему в онбординге, но разбор отложен."""
    topic = topic.strip()
    if not topic:
        return
    # Не дублируем
    existing = db_fetchone("SELECT id FROM parking_lot WHERE user_id=? AND topic=? AND discussed_at IS NULL",
                           (uid, topic))
    if existing:
        return
    db_exec("INSERT INTO parking_lot (user_id, topic, created_at) VALUES (?,?,?)",
            (uid, topic, datetime.now().isoformat()))

def get_parking_topics(uid: int, only_open: bool = True) -> list[tuple]:
    """Список тем. only_open=True — только неразобранные."""
    if only_open:
        return db_fetch("SELECT id, topic, created_at FROM parking_lot WHERE user_id=? AND discussed_at IS NULL ORDER BY id",
                        (uid,))
    return db_fetch("SELECT id, topic, created_at, discussed_at FROM parking_lot WHERE user_id=? ORDER BY id",
                    (uid,))

def mark_parking_discussed(topic_id: int):
    db_exec("UPDATE parking_lot SET discussed_at=? WHERE id=?",
            (datetime.now().isoformat(), topic_id))

# ── Подписки и лимиты ─────────────────────────────────────────────────────────

def ensure_subscription(uid: int):
    """Если у юзера ещё нет записи подписки — создаём триал на 7 дней.
    Вызывается при первом обращении. Повторные вызовы безопасны — INSERT OR IGNORE."""
    row = db_fetchone("SELECT user_id FROM subscriptions WHERE user_id=?", (uid,))
    if row:
        return
    cfg = PLANS["trial"]
    valid_until = (datetime.now() + timedelta(days=cfg["period_days"])).isoformat()
    db_exec("""INSERT OR IGNORE INTO subscriptions (user_id, plan, valid_until, last_payment_stars, last_payment_at)
               VALUES (?,?,?,?,?)""",
            (uid, "trial", valid_until, 0, datetime.now().isoformat()))

def get_user_plan(uid: int) -> str:
    """Возвращает текущий активный тариф: 'trial' | 'basic' | 'pro' | 'expired'.
    Если подписка истекла — возвращает 'expired' (не 'trial' — второго триала нет)."""
    row = db_fetchone("SELECT plan, valid_until FROM subscriptions WHERE user_id=?", (uid,))
    if not row:
        # Записи нет — значит это самый первый заход. Функция ensure_subscription()
        # создаст триал. Пока возвращаем 'trial' как дефолт.
        return "trial"
    plan, valid_until = row
    if not valid_until:
        return plan or "expired"
    try:
        if datetime.fromisoformat(valid_until) < datetime.now():
            return PLAN_EXPIRED
    except Exception:
        return PLAN_EXPIRED
    return plan

def activate_plan(uid: int, plan: str, stars_paid: int = 0):
    """Активировать/продлить тариф. Если подписка ещё активна — дни прибавляются к остатку."""
    cfg = PLANS.get(plan)
    if not cfg:
        return
    days = cfg["period_days"]
    now = datetime.now()
    if days > 0:
        current = db_fetchone("SELECT valid_until FROM subscriptions WHERE user_id=?", (uid,))
        base = now
        if current and current[0]:
            try:
                existing = datetime.fromisoformat(current[0])
                if existing > now:
                    base = existing
            except Exception:
                pass
        valid_until = (base + timedelta(days=days)).isoformat()
    else:
        valid_until = None
    db_exec("""INSERT INTO subscriptions (user_id, plan, valid_until, last_payment_stars, last_payment_at)
               VALUES (?,?,?,?,?)
               ON CONFLICT(user_id) DO UPDATE SET
                   plan=excluded.plan,
                   valid_until=excluded.valid_until,
                   last_payment_stars=excluded.last_payment_stars,
                   last_payment_at=excluded.last_payment_at""",
            (uid, plan, valid_until, stars_paid, now.isoformat()))

def days_left(uid: int) -> int:
    """Сколько дней осталось до конца текущей подписки. -1 если нет подписки."""
    row = db_fetchone("SELECT valid_until FROM subscriptions WHERE user_id=?", (uid,))
    if not row or not row[0]:
        return -1
    try:
        remaining = datetime.fromisoformat(row[0]) - datetime.now()
        return max(0, remaining.days)
    except Exception:
        return -1

def today_key(profile: dict | None = None) -> str:
    """Дата пользователя в формате YYYY-MM-DD — по его локальному времени.
    Если профиль не передан, по UTC (для внутренних счётчиков это не критично)."""
    if profile:
        return user_now(profile).date().isoformat()
    return datetime.now(timezone.utc).date().isoformat()

def get_usage(uid: int, kind: str, day: str | None = None) -> int:
    day = day or today_key()
    row = db_fetchone("SELECT count FROM usage_counters WHERE user_id=? AND day=? AND kind=?",
                      (uid, day, kind))
    return row[0] if row else 0

def bump_usage(uid: int, kind: str = "msg"):
    day = today_key()
    db_exec("""INSERT INTO usage_counters (user_id, day, kind, count)
               VALUES (?,?,?,1)
               ON CONFLICT(user_id, day, kind) DO UPDATE SET count = count + 1""",
            (uid, day, kind))

def check_plan_limit(uid: int, kind: str = "msg") -> tuple[bool, str]:
    """Проверяет доступ и дневной лимит.
    kind: 'msg' | 'voice' | 'photo'. Возвращает (ok, reason_text).
    Сам счётчик НЕ инкрементирует — это делает bump_usage() после успешной обработки."""
    # Если монетизация выключена флагом — всем всё разрешено, без лимитов.
    if not PAYMENTS_ENABLED:
        return True, ""
    plan = get_user_plan(uid)
    if plan == PLAN_EXPIRED:
        return False, (
            "⏸ Твой доступ к Нове закончился.\n\n"
            "Чтобы продолжить — оформи подписку: /subscribe\n"
            "Все твои задачи, цели и история сохранились — они подхватятся автоматически."
        )
    cfg = PLANS.get(plan, PLANS["trial"])
    limit_key = {"msg": "msg_daily", "voice": "voice_daily", "photo": "photo_daily"}.get(kind, "msg_daily")
    limit = cfg.get(limit_key, 0)
    used = get_usage(uid, kind)
    if used >= limit:
        upgrade_hint = "" if plan == "pro" else " Напиши /subscribe чтобы снять лимиты."
        kind_ru = {"msg": "сообщений", "voice": "голосовых", "photo": "фото"}.get(kind, kind)
        return False, (f"⏸ Ты достигла дневного лимита ({limit} {kind_ru}/день) на тарифе *{cfg['title']}*."
                       f"{upgrade_hint}")
    return True, ""

def user_has_feature(uid: int, feature: str) -> bool:
    """feature: 'calendar' | 'smart_model' | 'premium_ai' — проверка доступа по тарифу."""
    plan = get_user_plan(uid)
    if plan == PLAN_EXPIRED:
        return False
    cfg = PLANS.get(plan, PLANS["trial"])
    return bool(cfg.get(feature, False))

# ── Траты (распознанные чеки) ─────────────────────────────────────────────────

def add_expense(uid: int, amount: float, category: str, note: str = "", currency: str = "RUB"):
    db_exec("INSERT INTO expenses (user_id, amount, currency, category, note, created_at) VALUES (?,?,?,?,?,?)",
            (uid, amount, currency, category, note, datetime.now().isoformat()))

def get_expenses_summary(uid: int, days: int = 30) -> dict:
    """Сводка трат за N дней. Возвращает {'total': float, 'by_category': {cat: sum}, 'count': int}."""
    since = (datetime.now() - timedelta(days=days)).isoformat()
    rows = db_fetch("SELECT amount, category FROM expenses WHERE user_id=? AND created_at>=?",
                    (uid, since))
    total = sum(r[0] or 0 for r in rows)
    by_cat: dict = {}
    for amt, cat in rows:
        by_cat[cat or "прочее"] = by_cat.get(cat or "прочее", 0) + (amt or 0)
    return {"total": total, "by_category": by_cat, "count": len(rows)}

# ── Визуальные отчёты ─────────────────────────────────────────────────────────

def generate_sphere_chart(uid):
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        stats = get_sphere_stats(uid)
        goals = get_goals(uid)

        has_stats = bool(stats)
        has_goals = bool(goals)
        if not has_stats and not has_goals:
            return None

        fig, axes = plt.subplots(1, 2, figsize=(13, 5))
        fig.patch.set_facecolor('#1a1a2e')
        for ax in axes:
            ax.set_facecolor('#16213e')
            ax.tick_params(colors='#e0e0e0')
            ax.spines[:].set_color('#444')

        if has_stats:
            labels = [SPHERES.get(k, k).split()[-1] for k in stats.keys()]
            values = list(stats.values())
            bars = axes[0].barh(labels, values, color='#6C63FF', height=0.6)
            axes[0].set_title('Активность по сферам (7 дней)', color='#e0e0e0', pad=10)
            axes[0].set_xlabel('дней', color='#aaa')
            for bar, val in zip(bars, values):
                axes[0].text(bar.get_width() + 0.05, bar.get_y() + bar.get_height() / 2,
                             str(val), va='center', color='#e0e0e0', fontsize=9)
        else:
            axes[0].text(0.5, 0.5, 'Нет данных', ha='center', va='center',
                         color='#888', transform=axes[0].transAxes)
            axes[0].set_title('Активность по сферам', color='#e0e0e0')

        if has_goals:
            goal_names = [(g[1][:18] + '…' if len(g[1]) > 18 else g[1]) for g in goals[:6]]
            goal_vals = [g[4] for g in goals[:6]]
            colors = ['#FF6584' if v < 30 else '#FFCA3A' if v < 70 else '#6BCB77' for v in goal_vals]
            axes[1].barh(goal_names, goal_vals, color=colors, height=0.6)
            axes[1].set_xlim(0, 100)
            axes[1].set_title('Прогресс целей (%)', color='#e0e0e0', pad=10)
            axes[1].set_xlabel('%', color='#aaa')
            for i, val in enumerate(goal_vals):
                axes[1].text(val + 1, i, f'{val}%', va='center', color='#e0e0e0', fontsize=9)
        else:
            axes[1].text(0.5, 0.5, 'Целей нет', ha='center', va='center',
                         color='#888', transform=axes[1].transAxes)
            axes[1].set_title('Прогресс целей', color='#e0e0e0')

        plt.tight_layout(pad=2)
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logging.error(f"Chart error: {e}")
        return None

def generate_pdf_report(uid):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.lib import colors

        buf = io.BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, rightMargin=40, leftMargin=40,
                                topMargin=40, bottomMargin=40)
        styles = getSampleStyleSheet()
        h2 = ParagraphStyle('h2', parent=styles['Heading2'], textColor=colors.HexColor('#6C63FF'))
        normal = styles['Normal']
        story = []

        profile = get_profile(uid)
        name = profile.get("name", "Пользователь")
        now = datetime.now()

        story.append(Paragraph(f"Отчёт Nova — {name}", styles['Title']))
        story.append(Paragraph(now.strftime('%d.%m.%Y'), normal))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#6C63FF')))
        story.append(Spacer(1, 12))

        tasks = get_tasks(uid)
        story.append(Paragraph("📋 Открытые задачи", h2))
        if tasks:
            for t in tasks[:30]:
                icon = {'urgent': '🔴', 'important': '🟡'}.get(t[2], '⚪')
                story.append(Paragraph(f"{icon} {t[1]}", normal))
        else:
            story.append(Paragraph("Задач нет.", normal))
        story.append(Spacer(1, 12))

        goals = get_goals(uid)
        story.append(Paragraph("🎯 Активные цели", h2))
        if goals:
            for g in goals:
                progress = g[4] if len(g) > 4 else 0
                bar = '█' * (progress // 10) + '░' * (10 - progress // 10)
                story.append(Paragraph(f"• {g[1]}  [{bar}] {progress}%", normal))
        else:
            story.append(Paragraph("Целей нет.", normal))
        story.append(Spacer(1, 12))

        ideas = get_ideas(uid)
        story.append(Paragraph("💡 Идеи", h2))
        if ideas:
            for i in ideas[:20]:
                story.append(Paragraph(f"• {i[1]}", normal))
        else:
            story.append(Paragraph("Идей нет.", normal))

        doc.build(story)
        buf.seek(0)
        return buf
    except Exception as e:
        logging.error(f"PDF error: {e}")
        return None

# ── Mood / Energy / Habits / Journal / Wins ───────────────────────────────────

def log_mood(uid, score, note=""):
    db_exec("INSERT INTO mood_log (user_id,score,note,created_at) VALUES (?,?,?,?)",
            (uid, score, note, datetime.now().isoformat()))

def get_mood_history(uid, days=14):
    cutoff = (datetime.now() - timedelta(days=days)).isoformat()
    return db_fetch("SELECT score,note,created_at FROM mood_log WHERE user_id=? AND created_at>=? ORDER BY created_at",
                    (uid, cutoff))

def log_energy(uid, score):
    db_exec("INSERT INTO energy_log (user_id,score,created_at) VALUES (?,?,?)",
            (uid, score, datetime.now().isoformat()))

def get_energy_history(uid, days=14):
    cutoff = (datetime.now() - timedelta(days=days)).isoformat()
    return db_fetch("SELECT score,created_at FROM energy_log WHERE user_id=? AND created_at>=? ORDER BY created_at",
                    (uid, cutoff))

def get_or_create_habit(uid, name):
    existing = db_fetchone("SELECT id FROM habits WHERE user_id=? AND name=? AND active=1", (uid, name))
    if existing: return existing[0]
    db_exec("INSERT INTO habits (user_id,name,active,created_at) VALUES (?,?,1,?)",
            (uid, name, datetime.now().isoformat()))
    return db_fetchone("SELECT id FROM habits WHERE user_id=? AND name=? AND active=1", (uid, name))[0]

def get_habits(uid):
    return db_fetch("SELECT id,name FROM habits WHERE user_id=? AND active=1 ORDER BY id", (uid,))

def mark_habit_today(uid, habit_id):
    today = datetime.now().date().isoformat()
    existing = db_fetchone("SELECT id FROM habit_log WHERE user_id=? AND habit_id=? AND log_date=?",
                           (uid, habit_id, today))
    if not existing:
        db_exec("INSERT INTO habit_log (user_id,habit_id,log_date) VALUES (?,?,?)", (uid, habit_id, today))

def get_habit_streak(uid, habit_id):
    rows = db_fetch("SELECT log_date FROM habit_log WHERE user_id=? AND habit_id=? ORDER BY log_date DESC",
                    (uid, habit_id))
    if not rows: return 0
    streak = 0
    check = datetime.now().date()
    for (d,) in rows:
        if d == check.isoformat():
            streak += 1
            check -= timedelta(days=1)
        else:
            break
    return streak

def save_journal(uid, question, entry):
    db_exec("INSERT INTO journal (user_id,question,entry,created_at) VALUES (?,?,?,?)",
            (uid, question, entry, datetime.now().isoformat()))

def get_journal_entries(uid, limit=5):
    return db_fetch("SELECT question,entry,created_at FROM journal WHERE user_id=? ORDER BY created_at DESC LIMIT ?",
                    (uid, limit))

def add_win(uid, text):
    db_exec("INSERT INTO wins (user_id,text,created_at) VALUES (?,?,?)",
            (uid, text, datetime.now().isoformat()))

def get_wins(uid, limit=10):
    return db_fetch("SELECT text,created_at FROM wins WHERE user_id=? ORDER BY created_at DESC LIMIT ?",
                    (uid, limit))

# ── Колесо жизни ──────────────────────────────────────────────────────────────

def generate_wheel_chart(uid):
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        profile = get_profile(uid)
        spheres_score_str = profile.get("spheres_score", "")
        scores = {}
        if spheres_score_str:
            for pair in spheres_score_str.split(","):
                if ":" in pair:
                    k, v = pair.strip().split(":", 1)
                    try: scores[k.strip()] = int(v.strip())
                    except: pass

        labels_map = {
            "работа": "Работа", "финансы": "Финансы", "здоровье": "Здоровье",
            "отношения": "Отношения", "семья": "Семья", "саморазвитие": "Развитие",
            "творчество": "Творчество", "отдых": "Отдых",
            "духовность": "Духовность", "окружение": "Окружение",
        }
        keys = list(labels_map.keys())
        values = [scores.get(k, 5) for k in keys]
        labels = [labels_map[k] for k in keys]
        N = len(labels)
        angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
        values_plot = values + [values[0]]
        angles += angles[:1]

        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        fig.patch.set_facecolor('#1a1a2e')
        ax.set_facecolor('#16213e')
        ax.plot(angles, values_plot, 'o-', linewidth=2, color='#6C63FF')
        ax.fill(angles, values_plot, alpha=0.25, color='#6C63FF')
        ax.set_ylim(0, 10)
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(labels, color='#e0e0e0', fontsize=9)
        ax.tick_params(colors='#aaa')
        ax.yaxis.set_tick_params(labelcolor='#666')
        ax.spines['polar'].set_color('#444')
        ax.set_title('Колесо жизни', color='#e0e0e0', pad=15)

        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logging.error(f"Wheel chart error: {e}")
        return None

# ── Gantt / Mood-Energy / Habits charts ──────────────────────────────────────

def generate_gantt_chart(uid):
    """Диаграмма Ганта: задачи по срокам и приоритетам."""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        tasks = get_tasks(uid)
        if not tasks:
            return None

        today = datetime.now().date()
        PRIO_COLOR = {"urgent": "#FF6584", "important": "#FFCA3A", "normal": "#6C63FF"}
        TIMEFRAME_DAYS = {"today": 0, "week": 7, "month": 30, "quarter": 90, "year": 365}

        names, starts, durations, colors_list = [], [], [], []
        for t in tasks[:15]:
            tid, text, prio, sphere, tf, due = t
            label = (text[:28] + "…") if len(text) > 28 else text
            if due:
                try:
                    end = datetime.strptime(due[:10], "%Y-%m-%d").date()
                    days_left = max((end - today).days, 0)
                    dur = max(days_left, 1)
                    start_day = 0
                except Exception:
                    start_day = 0; dur = TIMEFRAME_DAYS.get(tf, 7)
            else:
                start_day = 0; dur = TIMEFRAME_DAYS.get(tf, 7)
            names.append(label); starts.append(start_day)
            durations.append(dur); colors_list.append(PRIO_COLOR.get(prio, "#6C63FF"))

        if not names:
            return None

        fig, ax = plt.subplots(figsize=(11, max(4, len(names) * 0.5 + 1)))
        fig.patch.set_facecolor('#1a1a2e')
        ax.set_facecolor('#16213e')

        y_pos = np.arange(len(names))
        ax.barh(y_pos, durations, left=starts, color=colors_list, height=0.6, alpha=0.85)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(names, color='#e0e0e0', fontsize=8)
        ax.tick_params(colors='#aaa')
        ax.spines[:].set_color('#444')
        ax.set_xlabel("Дней до дедлайна", color='#aaa', fontsize=9)
        ax.set_title("Диаграмма Ганта — задачи", color='#e0e0e0', pad=10)
        ax.axvline(x=0, color='#888', linestyle='--', linewidth=0.8)

        from matplotlib.patches import Patch
        legend = [Patch(color="#FF6584", label="Срочно"),
                  Patch(color="#FFCA3A", label="Важно"),
                  Patch(color="#6C63FF", label="Обычно")]
        ax.legend(handles=legend, facecolor='#1a1a2e', edgecolor='#444',
                  labelcolor='#e0e0e0', fontsize=8, loc='lower right')

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logging.error(f"Gantt chart error: {e}")
        return None


def generate_mood_energy_chart(uid):
    """График настроения и энергии за 14 дней."""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        mood_rows = get_mood_history(uid, 14)
        energy_rows = get_energy_history(uid, 14)
        if not mood_rows and not energy_rows:
            return None

        def rows_to_daily(rows, score_idx=0, date_idx=1):
            from collections import defaultdict
            daily = defaultdict(list)
            for row in rows:
                date_str = row[date_idx][:10]
                daily[date_str].append(row[score_idx])
            return {d: sum(v)/len(v) for d, v in daily.items()}

        mood_daily = rows_to_daily(mood_rows, score_idx=0, date_idx=2)
        energy_daily = rows_to_daily(energy_rows, score_idx=0, date_idx=1)

        all_dates = sorted(set(list(mood_daily.keys()) + list(energy_daily.keys())))
        if not all_dates:
            return None

        fig, ax = plt.subplots(figsize=(11, 4))
        fig.patch.set_facecolor('#1a1a2e')
        ax.set_facecolor('#16213e')
        ax.spines[:].set_color('#444')
        ax.tick_params(colors='#aaa')

        x_labels = [d[5:] for d in all_dates]  # MM-DD

        if mood_daily:
            mood_vals = [mood_daily.get(d) for d in all_dates]
            ax.plot(x_labels, mood_vals, 'o-', color='#FF6584', linewidth=2,
                    label='Настроение', markersize=5)
            ax.fill_between(x_labels, mood_vals, alpha=0.12, color='#FF6584')

        if energy_daily:
            energy_vals = [energy_daily.get(d) for d in all_dates]
            ax.plot(x_labels, energy_vals, 's-', color='#FFCA3A', linewidth=2,
                    label='Энергия', markersize=5)
            ax.fill_between(x_labels, energy_vals, alpha=0.12, color='#FFCA3A')

        ax.set_ylim(0, 11)
        ax.set_ylabel("Баллы (1-10)", color='#aaa', fontsize=9)
        ax.set_title("Настроение и энергия — 14 дней", color='#e0e0e0', pad=10)
        ax.legend(facecolor='#1a1a2e', edgecolor='#444', labelcolor='#e0e0e0', fontsize=9)
        plt.xticks(rotation=45, ha='right', fontsize=8)
        plt.tight_layout()

        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logging.error(f"Mood/energy chart error: {e}")
        return None


def generate_habit_chart(uid):
    """График выполнения привычек за 14 дней."""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        habits = get_habits(uid)
        if not habits:
            return None

        cutoff = (datetime.now() - timedelta(days=14)).date().isoformat()
        names, pcts = [], []
        for hid, hname in habits[:10]:
            rows = db_fetch(
                "SELECT COUNT(*) FROM habit_log WHERE user_id=? AND habit_id=? AND log_date>=?",
                (uid, hid, cutoff))
            done_days = rows[0][0] if rows else 0
            pct = round(done_days / 14 * 100)
            names.append((hname[:20] + "…") if len(hname) > 20 else hname)
            pcts.append(pct)

        if not names:
            return None

        fig, ax = plt.subplots(figsize=(9, max(3, len(names) * 0.55 + 1)))
        fig.patch.set_facecolor('#1a1a2e')
        ax.set_facecolor('#16213e')
        ax.spines[:].set_color('#444')
        ax.tick_params(colors='#aaa')

        colors_list = ['#6BCB77' if p >= 70 else '#FFCA3A' if p >= 40 else '#FF6584' for p in pcts]
        y_pos = np.arange(len(names))
        bars = ax.barh(y_pos, pcts, color=colors_list, height=0.6, alpha=0.85)
        ax.set_yticks(y_pos)
        ax.set_yticklabels(names, color='#e0e0e0', fontsize=9)
        ax.set_xlim(0, 105)
        ax.set_xlabel("Выполнение (%)", color='#aaa', fontsize=9)
        ax.set_title("Привычки — 14 дней", color='#e0e0e0', pad=10)
        for bar, pct in zip(bars, pcts):
            ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height() / 2,
                    f"{pct}%", va='center', color='#e0e0e0', fontsize=8)

        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=110, bbox_inches='tight')
        buf.seek(0)
        plt.close(fig)
        return buf
    except Exception as e:
        logging.error(f"Habit chart error: {e}")
        return None


# ── Инлайн-клавиатуры для трекеров ───────────────────────────────────────────


def habits_keyboard(uid):
    habits = get_habits(uid)
    today = datetime.now().date().isoformat()
    rows = []
    for hid, hname in habits:
        done = db_fetchone("SELECT id FROM habit_log WHERE user_id=? AND habit_id=? AND log_date=?",
                           (uid, hid, today))
        mark = "✅" if done else "⬜"
        streak = get_habit_streak(uid, hid)
        label = f"{mark} {hname}" + (f" 🔥{streak}" if streak > 1 else "")
        rows.append([InlineKeyboardButton(label, callback_data=f"habit_toggle_{hid}")])
    rows.append([InlineKeyboardButton("➕ Добавить привычку", callback_data="habit_add")])
    return InlineKeyboardMarkup(rows)

def settings_keyboard(profile):
    notif_morning = profile.get("notif_morning", "1") != "0"
    notif_evening = profile.get("notif_evening", "1") != "0"
    notif_weekly  = profile.get("notif_weekly",  "1") != "0"
    style = profile.get("info_style", "detailed")
    feedback = profile.get("feedback_style", "soft")
    return InlineKeyboardMarkup([
        [InlineKeyboardButton(f"{'🔔' if notif_morning else '🔕'} Утро", callback_data="set_morning"),
         InlineKeyboardButton(f"{'🔔' if notif_evening else '🔕'} Вечер", callback_data="set_evening"),
         InlineKeyboardButton(f"{'🔔' if notif_weekly  else '🔕'} Неделя", callback_data="set_weekly")],
        [InlineKeyboardButton(f"Стиль: {'кратко' if style=='brief' else 'подробно'}", callback_data="set_style"),
         InlineKeyboardButton(f"Тон: {'прямой' if feedback=='direct' else 'мягкий'}", callback_data="set_feedback")],
        [InlineKeyboardButton("🕐 Изменить часовой пояс", callback_data="set_tz")],
    ])

def add_goal(uid, text, sphere="general", timeframe="longterm"):
    existing = db_fetch("SELECT id FROM goals WHERE user_id=? AND text=? AND done=0", (uid, text))
    if existing: return
    db_exec("INSERT INTO goals (user_id,text,sphere,timeframe,created_at) VALUES (?,?,?,?,?)",
            (uid, text, sphere, timeframe, datetime.now().isoformat()))

def get_goals(uid, sphere=None, timeframe=None):
    query = "SELECT id,text,sphere,timeframe,progress FROM goals WHERE user_id=? AND done=0"
    params = [uid]
    if sphere: query += " AND sphere=?"; params.append(sphere)
    if timeframe: query += " AND timeframe=?"; params.append(timeframe)
    return db_fetch(query, tuple(params))

def add_idea(uid, text, sphere="general"):
    existing = db_fetch("SELECT id FROM ideas WHERE user_id=? AND text=?", (uid, text))
    if existing: return
    db_exec("INSERT INTO ideas (user_id,text,sphere,created_at) VALUES (?,?,?,?)",
            (uid, text, sphere, datetime.now().isoformat()))

def get_ideas(uid, sphere=None):
    if sphere:
        return db_fetch("SELECT id,text,sphere,created_at FROM ideas WHERE user_id=? AND sphere=?", (uid, sphere))
    return db_fetch("SELECT id,text,sphere,created_at FROM ideas WHERE user_id=?", (uid,))

def get_frozen_items(uid):
    week_ago = (datetime.now() - timedelta(days=7)).isoformat()[:10]
    ideas = db_fetch("SELECT id,text,sphere,'idea' FROM ideas WHERE user_id=? AND (reviewed_at IS NULL OR reviewed_at < ?)", (uid, week_ago))
    goals = db_fetch("SELECT id,text,sphere,'goal' FROM goals WHERE user_id=? AND done=0 AND progress=0 AND created_at < ?", (uid, week_ago))
    return ideas[:3] + goals[:2]

def log_sphere_activity(uid, sphere):
    db_exec("INSERT INTO sphere_activity (user_id,sphere,activity_date) VALUES (?,?,?)",
            (uid, sphere, datetime.now().date().isoformat()))

def get_sphere_stats(uid):
    rows = db_fetch("""SELECT sphere, COUNT(*) FROM sphere_activity
                 WHERE user_id=? AND activity_date >= date('now', '-7 days')
                 GROUP BY sphere ORDER BY COUNT(*) DESC""", (uid,))
    return {r[0]: r[1] for r in rows}

def save_google_token(uid, creds):
    expiry_str = creds.expiry.isoformat() if creds.expiry else None
    db_exec("""INSERT OR REPLACE INTO google_tokens
               (user_id, token, refresh_token, token_uri, client_id, client_secret, scopes, expiry)
               VALUES (?,?,?,?,?,?,?,?)""",
            (uid, creds.token, creds.refresh_token, creds.token_uri,
             creds.client_id, creds.client_secret, json.dumps(list(creds.scopes)), expiry_str))

def get_google_token(uid):
    try:
        # Пробуем с колонкой expiry (новая схема)
        row = db_fetchone("SELECT user_id,token,refresh_token,token_uri,client_id,client_secret,scopes,expiry FROM google_tokens WHERE user_id=?", (uid,))
        if not row: return None
        from datetime import datetime as dt
        expiry = dt.fromisoformat(row[7]) if (len(row) > 7 and row[7]) else None
        creds = Credentials(
            token=row[1], refresh_token=row[2], token_uri=row[3],
            client_id=row[4], client_secret=row[5],
            scopes=json.loads(row[6]), expiry=expiry)
        return creds
    except Exception:
        try:
            # Fallback без expiry (старая схема)
            row = db_fetchone("SELECT user_id,token,refresh_token,token_uri,client_id,client_secret,scopes FROM google_tokens WHERE user_id=?", (uid,))
            if not row: return None
            creds = Credentials(
                token=row[1], refresh_token=row[2], token_uri=row[3],
                client_id=row[4], client_secret=row[5],
                scopes=json.loads(row[6]))
            return creds
        except Exception as e:
            logging.error(f"get_google_token error uid={uid}: {e}")
            return None

async def get_calendar_service(uid):
    import asyncio
    creds = get_google_token(uid)
    if not creds: return None
    try:
        loop = asyncio.get_running_loop()
        # Обновляем токен если истёк или нет информации о сроке (на всякий случай)
        if creds.refresh_token and (creds.expired or creds.expiry is None):
            from google.auth.transport.requests import Request
            await loop.run_in_executor(None, lambda: creds.refresh(Request()))
            save_google_token(uid, creds)
        service = await loop.run_in_executor(
            None, lambda: build("calendar", "v3", credentials=creds,
                                cache_discovery=False))
        return service
    except Exception as e:
        logging.error(f"Calendar service error for uid={uid}: {e}")
        return None

async def add_to_calendar(uid, task_text, due_date=None, timeframe=None, event_time=None, priority="normal", description=""):
    import asyncio
    service = await get_calendar_service(uid)
    if not service:
        logging.warning(f"Calendar service unavailable for user {uid}")
        return None
    try:
        color_id = PRIORITY_COLOR.get(priority, "7")
        now = datetime.now()
        if due_date:
            start_date = due_date
        elif timeframe == "today":
            start_date = now.date().isoformat()
        elif timeframe == "tomorrow":
            start_date = (now + timedelta(days=1)).date().isoformat()
        elif timeframe == "week":
            start_date = (now + timedelta(days=3)).date().isoformat()
        elif timeframe == "month":
            start_date = (now + timedelta(days=14)).date().isoformat()
        else:
            start_date = (now + timedelta(days=7)).date().isoformat()

        months = ["янв","фев","мар","апр","май","июн","июл","авг","сен","окт","ноя","дек"]
        from datetime import date as date_type
        d = date_type.fromisoformat(start_date)
        date_label = f"{d.day} {months[d.month-1]}"

        if event_time:
            import pytz
            tz_str = get_profile(uid).get("timezone", "Europe/Moscow")
            try:
                tz = pytz.timezone(tz_str)
            except Exception:
                tz = pytz.timezone("Europe/Moscow")
            h, m = (event_time.split(":") + ["00"])[:2]
            start_dt = datetime.fromisoformat(start_date).replace(hour=int(h), minute=int(m))
            end_dt = start_dt + timedelta(hours=1)
            event = {
                "summary": task_text,
                "colorId": color_id,
                "description": description,
                "start": {"dateTime": tz.localize(start_dt).isoformat(), "timeZone": tz_str},
                "end":   {"dateTime": tz.localize(end_dt).isoformat(),   "timeZone": tz_str},
            }
            date_label = f"{d.day} {months[d.month-1]} в {h}:{m.zfill(2)}"
        else:
            event = {
                "summary": task_text,
                "colorId": color_id,
                "description": description,
                "start": {"date": start_date},
                "end":   {"date": start_date},
            }

        await asyncio.get_running_loop().run_in_executor(
            None, lambda: service.events().insert(calendarId="primary", body=event).execute()
        )
        logging.info(f"Calendar event added for user {uid}: {task_text} @ {start_date}")
        return date_label
    except Exception as e:
        logging.error(f"Calendar add error for user {uid}: {e}")
        return None

async def list_calendar_events(uid, max_results=30, include_past=False):
    """Возвращает список событий [{id, summary, start, calendar_id}]"""
    import asyncio
    service = await get_calendar_service(uid)
    if not service: return []
    try:
        loop = asyncio.get_running_loop()
        kwargs = dict(calendarId="primary", maxResults=max_results,
                      singleEvents=True, orderBy="startTime")
        if not include_past:
            kwargs["timeMin"] = datetime.utcnow().isoformat() + "Z"
        result = await loop.run_in_executor(None, lambda: service.events().list(**kwargs).execute())
        items = result.get("items", [])
        events = []
        for e in items:
            start = e.get("start", {})
            events.append({
                "id": e["id"],
                "summary": e.get("summary", "(без названия)"),
                "start": start.get("dateTime", start.get("date", "")),
            })
        return events
    except Exception as e:
        logging.error(f"Calendar list error uid={uid}: {e}")
        return []

async def delete_calendar_event(uid, event_id):
    import asyncio
    service = await get_calendar_service(uid)
    if not service: return False
    try:
        loop = asyncio.get_running_loop()
        await loop.run_in_executor(None, lambda: service.events().delete(
            calendarId="primary", eventId=event_id).execute())
        return True
    except Exception as e:
        logging.error(f"Calendar delete error uid={uid} event={event_id}: {e}")
        return False

async def update_calendar_event(uid, event_id, new_summary=None, new_date=None, new_time=None):
    import asyncio, pytz
    service = await get_calendar_service(uid)
    if not service: return False
    try:
        loop = asyncio.get_running_loop()
        event = await loop.run_in_executor(None, lambda: service.events().get(
            calendarId="primary", eventId=event_id).execute())
        if new_summary:
            event["summary"] = new_summary
        if new_date:
            tz_str = get_profile(uid).get("timezone", "Europe/Moscow")
            if new_time:
                try:
                    tz = pytz.timezone(tz_str)
                except Exception:
                    tz = pytz.timezone("Europe/Moscow")
                h, m = (new_time.split(":") + ["00"])[:2]
                start_dt = datetime.fromisoformat(new_date).replace(hour=int(h), minute=int(m))
                end_dt = start_dt + timedelta(hours=1)
                event["start"] = {"dateTime": tz.localize(start_dt).isoformat(), "timeZone": tz_str}
                event["end"]   = {"dateTime": tz.localize(end_dt).isoformat(),   "timeZone": tz_str}
            else:
                event["start"] = {"date": new_date}
                event["end"]   = {"date": new_date}
        await loop.run_in_executor(None, lambda: service.events().update(
            calendarId="primary", eventId=event_id, body=event).execute())
        return True
    except Exception as e:
        logging.error(f"Calendar update error uid={uid} event={event_id}: {e}")
        return False

async def list_calendars(uid):
    """Возвращает список всех календарей пользователя [{id, summary, primary}]"""
    import asyncio
    service = await get_calendar_service(uid)
    if not service: return []
    try:
        loop = asyncio.get_running_loop()
        result = await loop.run_in_executor(None, lambda: service.calendarList().list().execute())
        cals = []
        for c in result.get("items", []):
            cals.append({
                "id": c["id"],
                "summary": c.get("summary", "(без названия)"),
                "primary": c.get("primary", False),
            })
        return cals
    except Exception as e:
        logging.error(f"Calendar list error uid={uid}: {e}")
        return []

async def delete_extra_calendar(uid, calendar_id):
    """Удаляет дополнительный календарь (не primary)"""
    import asyncio
    if calendar_id == "primary":
        return False
    service = await get_calendar_service(uid)
    if not service: return False
    try:
        loop = asyncio.get_running_loop()
        # Пробуем удалить, если нет прав — просто скрываем из списка
        try:
            await loop.run_in_executor(None, lambda: service.calendars().delete(calendarId=calendar_id).execute())
        except Exception:
            await loop.run_in_executor(None, lambda: service.calendarList().delete(calendarId=calendar_id).execute())
        return True
    except Exception as e:
        logging.error(f"Calendar delete error uid={uid} cal={calendar_id}: {e}")
        return False

async def clear_all_calendar_events(uid):
    """Удаляет ВСЕ события (прошлые и будущие) из основного календаря"""
    events = await list_calendar_events(uid, max_results=500, include_past=True)
    deleted = 0
    for e in events:
        ok = await delete_calendar_event(uid, e["id"])
        if ok:
            deleted += 1
    return deleted

def get_oauth_flow():
    import os
    os.environ["OAUTHLIB_INSECURE_TRANSPORT"] = "1"
    client_config = {
        "web": {
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "auth_uri": "https://accounts.google.com/o/oauth2/auth",
            "token_uri": "https://oauth2.googleapis.com/token",
            "redirect_uris": [f"{WEBHOOK_URL}/oauth/callback"]
        }
    }
    flow = Flow.from_client_config(
        client_config,
        scopes=SCOPES,
        redirect_uri=f"{WEBHOOK_URL}/oauth/callback"
    )
    return flow

# Цвета Google Calendar по приоритету задач
PRIORITY_COLOR = {"urgent": "11", "important": "6", "normal": "7"}

def format_tasks(tasks, with_actions=False):
    if not tasks: return "Пусто 👌"
    icons = {"urgent": "🔴", "important": "🟡", "normal": "⚪"}
    lines = []
    for t in tasks:
        icon = icons.get(t[2], "⚪")
        lines.append(f"{icon} [{t[0]}] {t[1]}")
    return "\n".join(lines)

def format_goals(goals):
    if not goals: return "Пусто 👌"
    lines = []
    for g in goals:
        progress = g[4] if len(g) > 4 else 0
        prog_str = f" — {progress}%" if progress else ""
        lines.append(f"• [{g[0]}] {g[1]}{prog_str}")
    return "\n".join(lines)

def format_week_plan(uid):
    now = datetime.now()
    lines = [f"📅 *План на неделю*\n{now.strftime('%d.%m')} — {(now + timedelta(days=6)).strftime('%d.%m')}\n"]
    today = get_today_tasks(uid)
    week = get_tasks(uid, timeframe="week")
    urgent = get_tasks(uid, priority="urgent")
    if urgent:
        lines.append("🔴 *Срочно:*")
        for t in urgent[:5]: lines.append(f"• {t[1]}")
        lines.append("")
    if today:
        lines.append("📅 *Сегодня:*")
        for t in today[:5]: lines.append(f"• {t[1]}")
        lines.append("")
    if week:
        lines.append("📆 *На этой неделе:*")
        for t in week[:7]: lines.append(f"• {t[1]}")
    if not urgent and not today and not week:
        lines.append("Задач нет — отличное время добавить новые 🙂")
    return "\n".join(lines)

def format_dashboard(uid):
    profile = get_profile(uid)
    name = profile.get("name", "")
    today_tasks = get_today_tasks(uid)
    all_tasks = get_tasks(uid)
    goals = get_goals(uid)
    ideas = get_ideas(uid)
    stats = get_sphere_stats(uid)
    urgent = [t for t in all_tasks if t[2] == "urgent"]
    now = datetime.now()
    gcal = get_google_token(uid)
    cal_status = "✅ подключён" if gcal else "❌ не подключён"
    lines = [
        f"📊 {name} — {now.strftime('%d.%m.%Y')}",
        "",
        f"📅 Сегодня: {len(today_tasks)} задач",
        f"📌 Всего: {len(all_tasks)} | 🔴 Срочных: {len(urgent)}",
        f"🎯 Целей: {len(goals)} | 💡 Идей: {len(ideas)}",
        f"📆 Google Календарь: {cal_status}",
    ]
    if stats:
        lines.append("")
        lines.append("Активность за 7 дней:")
        inactive = set(SPHERE_KEYS) - set(stats.keys())
        for sk, cnt in stats.items():
            lines.append(f"  {SPHERES.get(sk, sk)}: {'▓' * min(cnt, 8)}")
        if inactive:
            lines.append("😴 Без внимания:")
            for s in list(inactive)[:2]:
                lines.append(f"  {SPHERES.get(s, s)}")
    return "\n".join(lines)

ONBOARDING_INTRO = """Привет! Я Нова — твой личный ассистент нового поколения 😎
Предлагаю на ты, но скажи как тебе комфортнее.
Я здесь чтобы помогать вести дела, двигаться к целям и развиваться 👍
Чем лучше тебя узнаю — тем точнее помогу.
Онбординг займёт время, но это инвестиция. Если не хочешь — скажи, сразу перейдём к делам.

Вот что нас ждёт:
📍 Этап 1 — Знакомство
📍 Этап 2 — Как ты устроен
📍 Этап 3 — Сферы жизни
📍 Этап 4 — Всё что есть прямо сейчас
📍 Этап 5 — Цели, мечты, идеи

Поехали — узнаем друг друга! 🚀"""

def build_system(profile, onboarding_mode=False, uid=None, cal_events=None):
    address = profile.get("address") or profile.get("name") or ""
    p_lines = []
    if address: p_lines.append(f"Обращение: {address}")
    if profile.get("occupation"): p_lines.append(f"Работа/деятельность: {profile['occupation']}")
    if profile.get("goals"): p_lines.append(f"Жизненные цели: {profile['goals']}")
    if profile.get("pain"): p_lines.append(f"Что мешает: {profile['pain']}")
    if profile.get("satisfied"): p_lines.append(f"Что хорошо в жизни: {profile['satisfied']}")
    if profile.get("day_rhythm"): p_lines.append(f"Ритм дня: {profile['day_rhythm']}")
    if profile.get("timezone"): p_lines.append(f"Часовой пояс: {profile['timezone']}")
    if profile.get("character"): p_lines.append(f"Характер/особенности: {profile['character']}")
    if profile.get("notif_extras"): p_lines.append(f"Доп. блоки в уведомлениях: {profile['notif_extras']}")
    if profile.get("notes"): p_lines.append(f"Заметки: {profile['notes']}")
    profile_block = "\n".join(p_lines)

    now = user_now(profile)
    current_time = f"Сейчас: {now.strftime('%A, %d.%m.%Y, %H:%M')}"

    # Блок с текущими событиями для управления календарём
    if cal_events:
        ev_lines = ["Текущие события в основном календаре:"]
        cals_list = None
        for e in cal_events:
            if "_calendars" in e:
                cals_list = e["_calendars"]
            else:
                ev_lines.append("• id=" + e['id'] + " | " + e['summary'] + " | " + e['start'])
        if cals_list:
            ev_lines.append("\nДополнительные календари (папки):")
            for c in cals_list:
                marker = "(основной)" if c.get("primary") else ""
                ev_lines.append("• id=" + c['id'] + " | " + c['summary'] + " " + marker)
        cal_events_block = "\n".join(ev_lines)
    else:
        cal_events_block = "(события не загружены)"

    # Блок Calendar для системного промпта
    if uid and get_google_token(uid):
        if onboarding_mode:
            cal_block = (
                "✅ Google Календарь подключён.\n\n"
                "ВАЖНО ВО ВРЕМЯ ОНБОРДИНГА: НЕ добавляй задачи в Google Calendar. "
                "НЕ упоминай что добавила в календарь. "
                "Задачи сохраняются в базе через [TASK:] теги. "
                "После онбординга предложишь запланировать всё вместе с датами и временем."
            )
        else:
            cal_block = (
                "✅ Google Календарь ПОДКЛЮЧЁН и РАБОТАЕТ.\n\n"
                "ВАЖНО: НИКОГДА не говори что у тебя нет доступа к календарю. "
                "НИКОГДА не предлагай подключить — уже подключён. "
                "Управляй через теги в конце ответа.\n\n"
                "ДОБАВИТЬ СРАЗУ → [TASK: название | приоритет | сфера | timeframe | HH:MM]\n"
                "УДАЛИТЬ ВСЕ → [CAL_DELETE_ALL]\n"
                "УДАЛИТЬ ОДНО → [CAL_DELETE: event_id]\n"
                "ИЗМЕНИТЬ → [CAL_UPDATE: event_id | название | YYYY-MM-DD | HH:MM]\n"
                "УДАЛИТЬ ПАПКУ → [CAL_DELETE_CALENDAR: calendar_id]\n\n"
                "ПЛАНИРОВАНИЕ С ПОДТВЕРЖДЕНИЕМ (когда просят расставить задачи по времени):\n"
                "1. Покажи список задач красиво\n"
                "2. Уточни для каждой: дата, время начала, сколько займёт (вероятная занятость)\n"
                "   Пример: 'Встреча в 14:00 — вероятно до 16:00, после можно ставить'\n"
                "3. Предложи итоговое расписание тегами:\n"
                "   [CAL_PLAN: название | YYYY-MM-DD | HH:MM | приоритет | заметки]\n"
                "4. Покажи красивую сводку:\n"
                "   🔴 Срочное — название — 15 апр 10:00 (~2ч, занят до 12:00)\n"
                "   🟡 Важное — название — 16 апр 14:00\n"
                "   🔵 Обычное — название — 17 апр\n"
                "   + описание задачи, заметки если есть\n"
                "5. Обязательно спроси: 'Всё подходит? Вношу в Google Calendar? ✅'\n"
                "6. ТОЛЬКО после 'да/вноси/ок' — добавляй. До этого — только [CAL_PLAN:] теги!\n\n"
                "ЦВЕТА: urgent=🔴(красный), important=🟡(оранжевый), normal=🔵(синий/голубой)\n"
                "Можно предлагать изменить цвет события — например 'Личные дела отмечу синим'\n\n"
                + cal_events_block
            )
    else:
        cal_block = "❌ Не подключён. Пользователь может подключить через /calendar. Не говори что у тебя нет интеграции с календарём — она есть, просто ещё не настроена."

    onboarding_block = ""
    if onboarding_mode:
        onboarding_block = """
═══════════════════════════════════
РЕЖИМ ОНБОРДИНГА — БЫСТРО И ПО ДЕЛУ
═══════════════════════════════════

ГЛАВНАЯ ЦЕЛЬ: за минимум сообщений собрать ключевые данные о человеке,
чтобы дальше работать с ним осмысленно. НЕ превращай онбординг в терапию
или длинный разговор — это быстрый старт. Глубокие темы поднимутся потом.

ТВОЙ ТЕМП:
- Один короткий вопрос → ответ → фиксируешь в [PROFILE:] → следующий
- Не задавай уточняющих вопросов в онбординге даже если очень хочется
- Если человек даёт развёрнутый ответ — коротко подтверди («поняла») и переходи дальше
- Не анализируй, не интерпретируй, не давай советов во время онбординга
- Перед каждым этапом объявляй его одной короткой строкой: "📍 Этап N из 5 — Название"
- После каждого этапа — в конце сообщения пиши тег [STAGE: N] где N — только что завершённый этап

ЕСЛИ ЧЕЛОВЕК ЗАТРОНУЛ СЛОЖНУЮ ТЕМУ:
- НЕ разбирай её прямо сейчас. Онбординг — не место для глубокого разбора.
- Запомни: в конце ответа добавь тег [PARKING: краткое описание темы]
- В ответе коротко скажи: «Это важно — обязательно вернёмся после знакомства»
- Пример: если говорит «у меня выгорание», ты фиксируешь [PARKING: выгорание]
  и пишешь: «Слышу. Об этом поговорим отдельно после знакомства. Пока двигаемся дальше.»

ЕСЛИ ЧЕЛОВЕК НЕ ХОЧЕТ ОНБОРДИНГ:
- Сразу пиши [PROFILE: onboarding_skipped=true] и говоришь что готова работать
- Предлагаешь нажать кнопку «Завершить знакомство» чтобы открыть меню

КНОПКИ: у пользователя всегда есть две кнопки под чатом:
- «✅ Завершить знакомство» — можно нажать в любой момент, онбординг прервётся, откроется меню
- «💡 Узнай меня больше» — если сам хочет рассказать больше после онбординга

───────────────────────────────────
ЭТАП 1 из 5 — КТО ТЫ
Цель: имя, обращение, род занятий, город, часовой пояс.
───────────────────────────────────
Вопросы строго по одному:
1. «Как тебя зовут и как мне к тебе обращаться?»
   → [PROFILE: name=...] [PROFILE: address=...]
2. «Чем занимаешься?»
   → [PROFILE: occupation=...]
3. «В каком городе живёшь? Это нужно чтобы утренние сообщения приходили в твоё местное время.»
   → ВСЕГДА сохраняй город отдельно: [PROFILE: city=...]
   → если человек назвал страну/регион, тоже сохрани: [PROFILE: location=...]
4. «Сколько сейчас времени у тебя? (напиши свои часы)»
   → вычисли часовой пояс → [PROFILE: timezone=UTC+N]
В конце сообщения этапа: [STAGE: 1]

───────────────────────────────────
ЭТАП 2 из 5 — КАК ТЕБЕ УДОБНО
Цель: как общаться именно с этим человеком.
───────────────────────────────────
Одним сообщением три коротких вопроса подряд:
- «Кратко или развёрнуто отвечать?»  → [PROFILE: info_style=brief/detailed]
- «Прямо или мягко давать обратную связь?»  → [PROFILE: feedback_style=direct/soft]
- «Что тебя заряжает — чтобы я знала и напоминала?»  → [PROFILE: energizers=...]
В конце: [STAGE: 2]

───────────────────────────────────
ЭТАП 3 из 5 — КОЛЕСО ЖИЗНИ
Цель: 10 сфер + приоритет.
───────────────────────────────────
Одним сообщением: список 10 сфер, просишь оценить каждую от 1 до 10
через запятую (например: «работа 7, финансы 5, здоровье 8...»).
Сферы: Работа, Финансы, Здоровье, Отношения, Семья, Саморазвитие, Творчество, Отдых, Духовность, Окружение.
После ответа: сохрани [PROFILE: spheres_score=работа:7,финансы:5,...]
Затем одним коротким вопросом: «Какая сфера в приоритете сейчас?»
→ [PROFILE: priority_sphere=...]
В конце: [STAGE: 3]

───────────────────────────────────
ЭТАП 4 из 5 — ВЫГРУЗКА
Цель: собрать текущие задачи и тревоги.
───────────────────────────────────
Одним сообщением: «Выгрузи всё что сейчас висит — дела, идеи, тревоги, планы.
Одним сообщением, списком или потоком — как удобно. Сама разберу.»
Из ответа парси:
- Срочные дела → [TASK: ... | urgent | general | today]
- Обычные дела → [TASK: ... | normal | general | week]
- Идеи и желания → [IDEA: ... | general]
- Если человек упомянул выгорание/тревогу/конфликт — [PARKING: тема]
Коротко подтверди: «Записала X задач и Y идей.»
В конце: [STAGE: 4]

───────────────────────────────────
ЭТАП 5 из 5 — КУДА ДВИЖЕМСЯ
Цель: большая цель + 2-3 краткосрочных + что мешает.
───────────────────────────────────
Три коротких вопроса по одному:
1. «Какая у тебя большая цель или мечта — на год-два вперёд?»
   → [GOAL: ... | priority_sphere | longterm]
2. «Назови 2-3 конкретные цели на ближайшие 3 месяца.»
   → [GOAL: ... | сфера | short] для каждой
3. «Что чаще всего мешает двигаться к целям?»
   → [PROFILE: pain=...]
В конце: [STAGE: 5]

───────────────────────────────────
ЗАВЕРШЕНИЕ
───────────────────────────────────
После [STAGE: 5]:
- Короткое (3-4 строки) тёплое резюме — 2-3 конкретных наблюдения о человеке
- Если в PARKING что-то копилось — напомни: «Отложили на потом: [темы]. Вернёмся когда скажешь.»
- Скажи что готова к работе и предложи нажать кнопку «✅ Завершить знакомство» для открытия меню
- НЕ пиши /done сама — кнопка это делает
"""

    return f"""Ты — Нова. Профессиональный личный ассистент.

{current_time}

ХАРАКТЕР И СТИЛЬ:
- Дружелюбная, тёплая, вдумчивая. Говоришь о себе "я", никогда не звучишь как бот.
- Не пишешь: "Конечно!", "Я понял!", "Как я могу помочь?" — только живые, человечные фразы.
- Обращаешься к человеку по имени/обращению из профиля когда уместно.
- Если человек пишет не на русском — отвечаешь на его языке.
- Замечаешь настроение и состояние — реагируешь с заботой.
- С каждым пользователем выстраиваешь персональные отношения на основе его профиля.

ГЛАВНАЯ ЗАДАЧА:
Разгружать голову. Принимать всё — задачи, идеи, планы — и превращать в чёткую структуру.
Фиксировать, распределять, отслеживать, напоминать, помогать двигаться вперёд.

ПРИОРИТЕТЫ:
1. Задачи и планирование — основное
2. Поддержка и советы
3. Сферы жизни — отслеживать баланс
4. Коучинг и рефлексия — только если человек просит

ФОРМАТИРОВАНИЕ (Telegram Markdown):
- *жирный* для важного, _курсив_ для акцентов
- Смайлики — уместно, не в каждой строке
- СТРОГО не более 4 строк в одном сообщении
- Без вступлений и предисловий — сразу по делу
- Списки через • когда нужно перечислить

ПРАВИЛО ПРО ЗАДАЧИ:
- "надо сделать", "запиши", "добавь" → фиксируй сам
- Человек рассуждает → спроси: "Добавить как задачу?"
- "хотелось бы", "мечтаю", "было бы здорово" → всегда фиксируй как идею без вопроса

УМНОЕ ПЛАНИРОВАНИЕ (обязательно):
- Если время задачи размыто ("на неделе", "как-нибудь", "скоро", "надо бы") — ВСЕГДА уточни конкретный день или предложи сам: "Поставить на вторник в 11:00?"
- Если у человека уже есть задачи — учти загруженность. Не ставь 5 срочных дел на один день.
- Для звонков, встреч, административных дел — предлагай конкретное время: утро (9-11), день (13-15), вечер (18-20).
- Если человек говорит "позвони", "запишись", "сходи" — это задача с временем. Спроси: "Когда удобнее — утром или днём?"
- После добавления задачи — кратко подтверди: что зафиксировала и когда.

РЕДАКТИРОВАНИЕ:
"удали задачу 5" → [DEL_TASK: 5]
"выполни задачу 3" → [DONE_TASK: 3]
"перенеси задачу 2 на неделю" → [EDIT_TASK: 2 | timeframe=week]
"измени задачу 1" → [EDIT_TASK: 1 | text=новый текст]
"прогресс по цели 4 — 60%" → [GOAL_PROGRESS: 4 | 60]

ЗАМОРОЖЕННЫЕ ЭЛЕМЕНТЫ:
Когда уместно — поднимай идеи и цели без движения. Предлагай запланировать.

ФИКСИРУЙ В КОНЦЕ ОТВЕТА (скрыто от пользователя):
[TASK: текст | приоритет | сфера | timeframe]
[GOAL: текст | сфера | timeframe]
[IDEA: текст | сфера]
[PROFILE: ключ=значение]
[DEL_TASK: id]
[DONE_TASK: id]
[EDIT_TASK: id | поле=значение]
[GOAL_PROGRESS: id | процент]

Приоритеты задач: urgent / important / normal
Timeframe: today / tomorrow / week / month / longterm
Время (5-е поле, опционально): формат HH:MM — указывай если пользователь назвал время или ты предлагаешь конкретный слот.
Пример с временем: [TASK: позвонить в ЖКХ | important | general | tomorrow | 11:00]
Пример без времени: [TASK: купить продукты | normal | general | today]
СФЕРЫ: {', '.join(SPHERES.values())}

GOOGLE CALENDAR:
{cal_block}

{onboarding_block}
{chr(10) + 'Профиль пользователя:' + chr(10) + profile_block if profile_block else ''}"""


async def process_response(uid, text, skip_calendar=False, cal_plan_buffer=None):
    cal_lines = []
    # Парсим все [TASK:...] теги ровно один раз — без дублирования
    for m in re.finditer(r'\[TASK:([^\]]+)\]', text):
        parts = [p.strip() for p in m.group(1).split('|')]
        date_str = None
        if len(parts) >= 5 and re.match(r'^\d{1,2}:\d{2}$', parts[4]):
            t, p, s, tf, et = parts[0], parts[1], parts[2], parts[3], parts[4]
            add_task(uid, t, p, s, tf)
            log_sphere_activity(uid, s)
            if not skip_calendar:
                date_str = await add_to_calendar(uid, t, timeframe=tf, event_time=et, priority=p)
        elif len(parts) >= 4:
            t, p, s, tf = parts[0], parts[1], parts[2], parts[3]
            add_task(uid, t, p, s, tf)
            log_sphere_activity(uid, s)
            if not skip_calendar:
                date_str = await add_to_calendar(uid, t, timeframe=tf, priority=p)
        elif len(parts) >= 3:
            t, p, s = parts[0], parts[1], parts[2]
            add_task(uid, t, p, s)
            if not skip_calendar:
                date_str = await add_to_calendar(uid, t, priority=p)
        else:
            continue
        if date_str:
            cal_lines.append(f"📆 Добавила в календарь: _{t}_ — {date_str}")

    # CAL_PLAN теги — собираем план для подтверждения (не добавляем сразу)
    if cal_plan_buffer is not None:
        for m in re.finditer(r'\[CAL_PLAN:([^\]]+)\]', text):
            parts = [p.strip() for p in m.group(1).split('|')]
            if len(parts) >= 2:
                cal_plan_buffer.append(parts)
    for t, s, tf in re.findall(r'\[GOAL:\s*(.+?)\s*\|\s*(\w+)\s*\|\s*(\w+)\s*\]', text):
        add_goal(uid, t, s, tf)
    for t, s in re.findall(r'\[GOAL:\s*([^|]+?)\s*\|\s*(\w+)\s*\]', text):
        add_goal(uid, t, s)
    for gid, val in re.findall(r'\[GOAL_PROGRESS:\s*(\d+)\s*\|\s*(\d+)\s*\]', text):
        update_goal_progress(int(gid), int(val))
    for t, s in re.findall(r'\[IDEA:\s*(.+?)\s*\|\s*(\w+)\s*\]', text):
        add_idea(uid, t, s)
    for tid in re.findall(r'\[DONE_TASK:\s*(\d+)\s*\]', text):
        complete_task(int(tid))
    for tid in re.findall(r'\[DEL_TASK:\s*(\d+)\s*\]', text):
        delete_task(int(tid))
    for match in re.findall(r'\[EDIT_TASK:\s*(\d+)\s*\|\s*(.+?)\s*\]', text):
        tid = int(match[0])
        for pair in match[1].split('|'):
            if '=' in pair:
                k, _, v = pair.partition('=')
                edit_task(tid, **{k.strip(): v.strip()})
    profile_matches = re.findall(r'\[PROFILE:\s*(.+?)\s*\]', text)
    if profile_matches:
        profile = get_profile(uid)
        for m in profile_matches:
            for pair in m.split(','):
                if '=' in pair:
                    k, _, v = pair.partition('=')
                    profile[k.strip()] = v.strip()
        save_profile(uid, profile)
    # Прогресс онбординга: Нова объявляет завершение этапа тегом [STAGE: N]
    for stage in re.findall(r'\[STAGE:\s*(\d+)\s*\]', text):
        try:
            update_user(uid, onboarding_step=int(stage))
        except Exception:
            pass
    # Отложенные темы (на время онбординга): [PARKING: краткое описание темы]
    for topic in re.findall(r'\[PARKING:\s*([^\]]+?)\s*\]', text):
        add_parking(uid, topic)
    # Управление событиями Google Calendar
    if re.search(r'\[CAL_DELETE_ALL\]', text):
        deleted = await clear_all_calendar_events(uid)
        if deleted > 0:
            cal_lines.append(f"🗑 Удалила {deleted} событий из календаря")
        else:
            cal_lines.append("🗑 События не найдены или Calendar не подключён")
    for event_id in re.findall(r'\[CAL_DELETE:\s*([^\]]+?)\s*\]', text):
        ok = await delete_calendar_event(uid, event_id.strip())
        if ok:
            cal_lines.append("🗑 Событие удалено из календаря")
    for cal_id in re.findall(r'\[CAL_DELETE_CALENDAR:\s*([^\]]+?)\s*\]', text):
        ok = await delete_extra_calendar(uid, cal_id.strip())
        if ok:
            cal_lines.append("🗑 Календарь удалён")
    for match in re.findall(r'\[CAL_UPDATE:\s*([^|]+?)\s*\|\s*([^|]*?)\s*\|\s*([^|]*?)\s*\|\s*([^\]]*?)\s*\]', text):
        event_id, new_title, new_date, new_time = match
        ok = await update_calendar_event(uid, event_id.strip(),
            new_summary=new_title.strip() or None,
            new_date=new_date.strip() or None,
            new_time=new_time.strip() or None)
        if ok:
            cal_lines.append(f"✏️ Событие обновлено в календаре")

    _TAGS = r'TASK|GOAL|IDEA|PROFILE|DONE_TASK|DEL_TASK|EDIT_TASK|GOAL_PROGRESS|CAL_DELETE_ALL|CAL_DELETE|CAL_UPDATE|CAL_DELETE_CALENDAR|CAL_PLAN|STAGE|PARKING|EXPENSE'
    # Полные теги (с закрывающей скобкой)
    text = re.sub(rf'\[({_TAGS}):[^\]]*\]', '', text, flags=re.DOTALL)
    text = re.sub(r'\[CAL_DELETE_ALL\]', '', text)
    # Обрезанные теги в конце сообщения (без закрывающей скобки — ответ оборвался)
    text = re.sub(rf'\[({_TAGS}):[^\]]*$', '', text.strip(), flags=re.DOTALL)
    result = text.strip()
    if cal_lines:
        result = result + "\n\n" + "\n".join(cal_lines)
    return result

async def send_safe(update, text, reply_markup=None):
    try:
        await update.message.reply_text(text, parse_mode="Markdown", reply_markup=reply_markup)
    except Exception:
        await update.message.reply_text(text, reply_markup=reply_markup)

async def handle_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    uid = query.from_user.id
    data = query.data

    async def edit(text, kb=None):
        try:
            await query.edit_message_text(text, parse_mode="Markdown", reply_markup=kb)
        except:
            await query.edit_message_text(text, reply_markup=kb)

    if data == "onboarding_start":
        update_user(uid, onboarding_step=1)
        # Деактивируем кнопку, оставляя приветственный текст нетронутым
        disabled_kb = InlineKeyboardMarkup([[
            InlineKeyboardButton("✅ Начато", callback_data="noop")
        ]])
        try:
            await query.edit_message_reply_markup(reply_markup=disabled_kb)
        except Exception:
            pass
        profile = get_profile(uid)
        system = build_system(profile, onboarding_mode=True, uid=uid)
        try:
            response = await call_claude(
                get_history(uid) + [{"role": "user", "content": "Начинаем! Старт этапа 1."}],
                system, model=MODEL_SMART)
            clean = await process_response(uid, response, skip_calendar=True)
            save_msg(uid, "user", "Начинаем!")
            save_msg(uid, "assistant", clean)
            await query.message.reply_text(clean, parse_mode="Markdown")
        except Exception as e:
            logging.error(f"Onboarding start error: {e}")
            await query.message.reply_text("Отлично, начнём! Как тебя зовут и как к тебе обращаться?")
        return
    if data == "noop":
        return
    # ── Кнопки онбординга (отображаются под каждым сообщением Новы) ──
    if data == "finish_onboarding":
        # Завершаем онбординг принудительно — как будто юзер написал /done
        update_user(uid, onboarding_done=1)
        clear_followup(uid)
        # Короткое тёплое завершение + основное меню
        try:
            await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            pass
        # Парковочные темы — напомним что они есть
        parked = get_parking_topics(uid, only_open=True)
        parked_hint = ""
        if parked:
            topics = ", ".join([t[1] for t in parked[:5]])
            parked_hint = f"\n\nКстати, мы отложили на потом: _{topics}_. Скажи когда захочешь вернуться — запиши себе или просто напиши мне."
        await query.message.reply_text(
            f"Принято. Открываю меню — дальше работаем.{parked_hint}",
            parse_mode="Markdown",
            reply_markup=main_keyboard())
        return
    if data == "know_me_more":
        # Юзер хочет рассказать больше — подсказываем ему, что можно написать
        try:
            await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            pass
        await query.message.reply_text(
            "Слушаю — расскажи что считаешь важным. Это может быть:\n"
            "• что у тебя на душе\n"
            "• чего ты хочешь\n"
            "• чего избегаешь\n"
            "• что беспокоит\n\n"
            "Пиши как есть, я запишу.",
            reply_markup=onboarding_keyboard())
        return
    # ── Покупка подписки ──
    if data in ("buy_basic", "buy_pro"):
        if not PAYMENTS_ENABLED:
            await query.message.reply_text("Оплата сейчас отключена. Попробуй позже.")
            return
        plan_key = data.split("_", 1)[1]
        try:
            await send_stars_invoice(context, uid, plan_key)
        except Exception as e:
            logging.error(f"send_invoice error uid={uid}: {e}")
            await query.message.reply_text(
                "Не удалось создать счёт. Убедись что у тебя последняя версия Telegram.")
        return
    # ── Забыть ──
    if data == "forget_chat":
        clear_history(uid)
        await edit("🧹 Короткая история диалога очищена. Профиль и задачи на месте.")
        return
    if data == "forget_memory":
        ok = mem0_delete_all_user(uid)
        clear_history(uid)
        await edit("🌫 Готово. Я забыла то что знала о тебе из прошлых разговоров — "
                   "профиль, задачи и цели остались.")
        return
    # ── NPS-ответ (цифра 1-10) ──
    if data.startswith("nps_"):
        try:
            rating = int(data.split("_", 1)[1])
        except Exception:
            return
        db_exec("""UPDATE feedback SET rating=?, answered_at=? WHERE user_id=?""",
                (rating, datetime.now().isoformat(), uid))
        if rating <= 6:
            await edit(f"Спасибо 🙏 Оценка *{rating}*. Что можно улучшить? Напиши мне прямо в чат — любую критику приму.",
                       None)
        elif rating <= 8:
            await edit(f"Оценка *{rating}* — принято 💛 Если есть что предложить — напиши, я учту.")
        else:
            await edit(f"Оценка *{rating}* 🔥 Очень рада. Поделись со мной тем что нравится больше всего — или кинь "
                       f"другу ссылку /invite.")
        return
    # ── Возврат к отложенной теме ──
    if data.startswith("discuss_"):
        try:
            tid = int(data.split("_", 1)[1])
        except Exception:
            return
        row = db_fetchone("SELECT topic FROM parking_lot WHERE id=? AND user_id=?", (tid, uid))
        if not row:
            await edit("Тема не найдена — возможно, уже разобрали.")
            return
        topic = row[0]
        mark_parking_discussed(tid)
        profile = get_profile(uid)
        system = build_system(profile, uid=uid)
        try:
            response = await call_claude(
                get_history(uid, limit=6) + [{"role": "user",
                    "content": f"Давай сейчас вернёмся к теме, которую мы отложили в онбординге: «{topic}». "
                               f"Задай один открытый вопрос, чтобы я начала рассказывать. Слушай и помогай распутать."}],
                system, model=MODEL_SMART, max_tokens=MAX_TOKENS_DEFAULT)
            clean = await process_response(uid, response)
            await query.message.reply_text(clean, parse_mode="Markdown", reply_markup=main_keyboard())
        except Exception as e:
            logging.error(f"parking discuss error: {e}")
            await query.message.reply_text(f"Окей, говорим про «{topic}». Расскажи что чувствуешь.",
                                           reply_markup=main_keyboard())
        return
    if data == "parking_clear":
        db_exec("UPDATE parking_lot SET discussed_at=? WHERE user_id=? AND discussed_at IS NULL",
                (datetime.now().isoformat(), uid))
        await edit("Хорошо — темы убрала из списка. Если всплывут снова, поговорим.")
        return
    # ── Подтверждение удаления данных (GDPR) ──
    if data == "delete_me_yes":
        wipe_user_data(uid)
        await edit("🗑 Все твои данные удалены. Было приятно работать вместе. Если передумаешь — /start.")
        return
    if data == "delete_me_no":
        await edit("Отменила удаление. Продолжаем 🌿")
        return
    if data == "back_main":
        await edit("Главное меню 👇"); return
    if data == "back_spheres":
        await edit("Выбери сферу:", spheres_keyboard()); return
    if data == "tasks_today":
        await edit(f"📅 *На сегодня:*\n\n{format_tasks(get_today_tasks(uid))}", tasks_keyboard()); return
    if data == "tasks_week":
        await edit(f"📆 *На неделю:*\n\n{format_tasks(get_tasks(uid, timeframe='week'))}", tasks_keyboard()); return
    if data == "tasks_month":
        await edit(f"🗓 *На месяц:*\n\n{format_tasks(get_tasks(uid, timeframe='month'))}", tasks_keyboard()); return
    if data == "tasks_longterm":
        await edit(f"♾ *Долгосрочные:*\n\n{format_tasks(get_tasks(uid, timeframe='longterm'))}", tasks_keyboard()); return
    if data == "tasks_urgent":
        await edit(f"🔴 *Срочные:*\n\n{format_tasks(get_tasks(uid, priority='urgent'))}", tasks_keyboard()); return
    if data == "tasks_done":
        await edit(f"✅ *Выполненные:*\n\n{format_tasks(get_tasks(uid, done=1))}", tasks_keyboard()); return
    if data == "tasks_all":
        tasks = get_tasks(uid)
        await edit(f"📋 *Все задачи:*\n\n{format_tasks(tasks)}", tasks_keyboard()); return
    if data == "goals_short":
        goals = get_goals(uid, timeframe="short")
        await edit(f"*⚡ Краткосрочные:*\n\n{format_goals(goals)}", goals_keyboard()); return
    if data == "goals_long":
        goals = get_goals(uid, timeframe="longterm")
        await edit(f"*🏔 Долгосрочные:*\n\n{format_goals(goals)}", goals_keyboard()); return
    if data == "goals_all":
        goals = get_goals(uid)
        await edit(f"*🎯 Все цели:*\n\n{format_goals(goals)}", goals_keyboard()); return
    if data.startswith("tdone_"):
        complete_task(int(data.replace("tdone_", "")))
        await edit("✅ Задача выполнена!", tasks_keyboard()); return
    if data.startswith("tdel_"):
        delete_task(int(data.replace("tdel_", "")))
        await edit("🗑 Задача удалена.", tasks_keyboard()); return
    if data.startswith("tmove_"):
        tid = data.replace("tmove_", "")
        await edit(f"Перенести задачу [{tid}] на:", move_timeframe_keyboard(tid)); return
    if data.startswith("tset_"):
        parts = data.split("_", 2)
        tf, tid = parts[1], int(parts[2])
        today = datetime.now().date().isoformat()
        tomorrow = (datetime.now() + timedelta(days=1)).date().isoformat()
        if tf == "today":
            edit_task(tid, timeframe="today", due_date=today)
        elif tf == "tomorrow":
            edit_task(tid, timeframe="today", due_date=tomorrow)
        elif tf == "week":
            edit_task(tid, timeframe="week", due_date=None)
        elif tf == "month":
            edit_task(tid, timeframe="month", due_date=None)
        await edit("✅ Задача перенесена!", tasks_keyboard()); return
    if data.startswith("sphere_"):
        sk = data.replace("sphere_", "")
        log_sphere_activity(uid, sk)
        await edit(f"{SPHERES.get(sk)}\n\nЧто смотрим?", sphere_detail_keyboard(sk)); return
    if data.startswith("sph_tasks_"):
        sk = data.replace("sph_tasks_", "")
        await edit(f"{SPHERES.get(sk)} — задачи:\n\n{format_tasks(get_tasks(uid, sphere=sk))}", sphere_detail_keyboard(sk)); return
    if data.startswith("sph_goals_"):
        sk = data.replace("sph_goals_", "")
        goals = get_goals(uid, sphere=sk)
        await edit(f"{SPHERES.get(sk)} — цели:\n\n{format_goals(goals)}", sphere_detail_keyboard(sk)); return
    if data.startswith("sph_ideas_"):
        sk = data.replace("sph_ideas_", "")
        ideas = get_ideas(uid, sphere=sk)
        text = f"{SPHERES.get(sk)} — идеи:\n\n" + ("\n".join([f"• {i[1]}" for i in ideas]) if ideas else "Пусто 👌")
        await edit(text, sphere_detail_keyboard(sk)); return
    if data.startswith("mood_"):
        score = int(data.split("_")[1])
        log_mood(uid, score)
        profile = get_profile(uid)
        hist = get_mood_history(uid, days=7)
        avg = round(sum(r[0] for r in hist) / len(hist), 1) if hist else score
        comment = ""
        if len(hist) >= 3:
            trend = hist[-1][0] - hist[0][0]
            if trend > 2: comment = " Заметный рост за неделю 📈"
            elif trend < -2: comment = " Сдаёт немного — следи за собой 💛"
        await edit(f"Записала настроение: *{score}/10*{comment}\nСредний за неделю: {avg}/10")
        return
    if data.startswith("energy_"):
        score = int(data.split("_")[1])
        log_energy(uid, score)
        hist = get_energy_history(uid, days=7)
        avg = round(sum(r[0] for r in hist) / len(hist), 1) if hist else score
        await edit(f"Записала энергию: *{score}/10*\nСредняя за неделю: {avg}/10")
        return
    if data.startswith("habit_toggle_"):
        hid = int(data.replace("habit_toggle_", ""))
        mark_habit_today(uid, hid)
        await query.edit_message_reply_markup(reply_markup=habits_keyboard(uid))
        return
    if data == "habit_add":
        await edit("Напиши название привычки которую хочешь добавить:")
        context.user_data["mode"] = "habit_add"
        return
    if data in ("set_morning", "set_evening", "set_weekly"):
        key = {"set_morning": "notif_morning", "set_evening": "notif_evening", "set_weekly": "notif_weekly"}[data]
        profile = get_profile(uid)
        current = profile.get(key, "1")
        profile[key] = "0" if current != "0" else "1"
        save_profile(uid, profile)
        await query.edit_message_reply_markup(reply_markup=settings_keyboard(profile))
        return
    if data == "set_style":
        profile = get_profile(uid)
        profile["info_style"] = "brief" if profile.get("info_style", "detailed") == "detailed" else "detailed"
        save_profile(uid, profile)
        await query.edit_message_reply_markup(reply_markup=settings_keyboard(profile))
        return
    if data == "set_feedback":
        profile = get_profile(uid)
        profile["feedback_style"] = "direct" if profile.get("feedback_style", "soft") == "soft" else "soft"
        save_profile(uid, profile)
        await query.edit_message_reply_markup(reply_markup=settings_keyboard(profile))
        return
    if data == "set_tz":
        await edit("Напиши своё текущее время — например '15:30' или '9 утра'. Я вычислю часовой пояс.")
        context.user_data["mode"] = "set_tz"
        return

async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)
    # Разбор реферального аргумента: /start ref_<id>
    args = context.args if hasattr(context, "args") else []
    if args and args[0].startswith("ref_"):
        try:
            inviter = int(args[0][4:])
            record_referral(uid, inviter)
            logging.info(f"Referral recorded: {inviter} invited {uid}")
        except Exception as e:
            logging.warning(f"Bad ref arg: {e}")
    user = get_user(uid)
    if user[1]:
        profile = get_profile(uid)
        name = profile.get("name", "")
        await send_safe(update, f"Я здесь, {name} 👋" if name else "Я здесь 👋", main_keyboard())
    else:
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🚀 Поехали!", callback_data="onboarding_start")]])
        await update.message.reply_text(ONBOARDING_INTRO, reply_markup=kb)
        save_msg(uid, "assistant", ONBOARDING_INTRO)

async def cmd_done(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    update_user(uid, onboarding_done=1)
    # После завершения онбординга — убираем незакрытые follow-up,
    # чтобы бот не возвращался к уже неактуальным вопросам настройки.
    clear_followup(uid)
    profile = get_profile(uid)
    system = build_system(profile, uid=uid)
    try:
        response = await call_claude(
            get_history(uid) + [{"role": "user", "content": "Знакомство завершено. Сделай краткий вывод — что знаешь обо мне и с чего начнём. Открой меню."}],
            system, model=MODEL_SMART)
        clean = await process_response(uid, response)
        save_msg(uid, "assistant", clean)
        await send_safe(update, clean, main_keyboard())
    except:
        await send_safe(update, "Отлично, поехали! 🚀", main_keyboard())

    # Показываем список задач и предлагаем запланировать в календарь
    tasks = get_tasks(uid)
    if not tasks:
        return
    icons = {"urgent": "🔴", "important": "🟡", "normal": "⚪"}
    task_lines = [f"{icons.get(t[2],'⚪')} {t[1]}" for t in tasks[:12]]
    extra = f"\n_...и ещё {len(tasks)-12}_" if len(tasks) > 12 else ""
    task_block = "\n".join(task_lines) + extra

    if get_google_token(uid):
        await update.message.reply_text(
            f"📋 У тебя уже *{len(tasks)} задач* в списке:\n\n{task_block}\n\n"
            f"Давай расставим их по времени в Google Calendar! "
            f"Напиши — и я спрошу про каждую: когда, во сколько, сколько займёт.",
            parse_mode="Markdown", reply_markup=main_keyboard())
    else:
        await update.message.reply_text(
            f"📋 Зафиксировала *{len(tasks)} задач*:\n\n{task_block}\n\n"
            f"Подключи Google Calendar (/calendar) — и я расставлю всё по времени!",
            parse_mode="Markdown", reply_markup=main_keyboard())

async def cmd_calendar(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    existing = get_google_token(uid)
    if existing:
        # Проверяем что токен реально работает
        service = await get_calendar_service(uid)
        if service:
            await update.message.reply_text(
                "✅ Google Календарь подключён и работает!\n\nВсе новые задачи автоматически попадают в календарь.\n\nЧтобы переподключить — напиши /calreset",
                reply_markup=main_keyboard())
            return
        else:
            # Токен есть но сломан — удаляем и переподключаем
            db_exec("DELETE FROM google_tokens WHERE user_id=?", (uid,))
            await update.message.reply_text(
                "⚠️ Токен Calendar устарел, нужно переподключить.\nСейчас отправлю новую ссылку...",
                reply_markup=main_keyboard())
    try:
        flow = get_oauth_flow()
        auth_url, _ = flow.authorization_url(
            access_type="offline",
            include_granted_scopes="true",
            state=str(uid),
            prompt="consent"
        )
        await update.message.reply_text(
            f"📆 Подключаем Google Календарь!\n\nНажми на ссылку, войди в Google и разреши доступ:\n\n{auth_url}\n\nПосле авторизации бот автоматически получит доступ.",
            reply_markup=main_keyboard()
        )
    except Exception as e:
        logging.error(f"Calendar auth error: {e}")
        await update.message.reply_text("Что-то пошло не так при подключении календаря(", reply_markup=main_keyboard())

async def cmd_calreset(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    db_exec("DELETE FROM google_tokens WHERE user_id=?", (uid,))
    await update.message.reply_text("🔄 Отключила Calendar. Напиши /calendar чтобы подключить заново.", reply_markup=main_keyboard())

async def cmd_calshow(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    if not get_google_token(uid):
        await update.message.reply_text("❌ Google Calendar не подключён. Используй /calendar", reply_markup=main_keyboard())
        return
    await update.message.reply_text("Загружаю события... 📅")
    try:
        events = await list_calendar_events(uid, max_results=30, include_past=False)
        cals   = await list_calendars(uid)
        lines  = ["📅 *Твой Google Calendar*\n"]
        if len(cals) > 1:
            lines.append(f"📁 *Календари ({len(cals)}):*")
            for c in cals:
                marker = " _(основной)_" if c.get("primary") else ""
                lines.append(f"  • {c['summary']}{marker}")
            lines.append("")
        if events:
            lines.append(f"*Предстоящих событий: {len(events)}*\n")
            for e in events[:20]:
                start = e['start'][:10] if e['start'] else "?"
                lines.append(f"• {e['summary']} — {start}")
            if len(events) > 20:
                lines.append(f"\n_...и ещё {len(events)-20}_")
            lines.append("\nНапиши что хочешь сделать — изучить, удалить лишнее, добавить новое.")
        else:
            lines.append("Событий нет. Хочешь запланировать задачи в календарь?")
        await send_safe(update, "\n".join(lines), main_keyboard())
    except Exception as e:
        logging.error(f"calshow error: {e}")
        await update.message.reply_text("Не смогла загрузить календарь( Попробуй /calreset и переподключи.", reply_markup=main_keyboard())

async def cmd_calinfo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    try:
        lines = ["Диагностика Google Calendar:\n"]
        lines.append("GOOGLE_CLIENT_ID: " + ("есть" if GOOGLE_CLIENT_ID else "НЕТ - добавить в Railway Variables"))
        lines.append("GOOGLE_CLIENT_SECRET: " + ("есть" if GOOGLE_CLIENT_SECRET else "НЕТ - добавить в Railway Variables"))
        lines.append("WEBHOOK_URL: " + str(WEBHOOK_URL))
        try:
            token = get_google_token(uid)
            lines.append("Токен в базе: " + ("есть" if token else "нет - нужно /calendar"))
        except Exception as e:
            lines.append("Токен в базе: ошибка - " + str(e)[:100])
            token = None
        if token:
            try:
                service = await get_calendar_service(uid)
                lines.append("Токен рабочий: " + ("да" if service else "нет - нужно /calreset потом /calendar"))
                expired = getattr(token, 'expired', False)
                lines.append("Токен истёк: " + ("да" if expired else "нет"))
            except Exception as e:
                lines.append("Проверка токена: ошибка - " + str(e)[:100])
        await update.message.reply_text("\n".join(lines), reply_markup=main_keyboard())
    except Exception as e:
        await update.message.reply_text("Ошибка диагностики: " + str(e)[:200])

async def cmd_newuser(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    wipe_user_data(uid)
    await update.message.reply_text("Сброс выполнен. Напиши /start")

async def cmd_reset(update: Update, context: ContextTypes.DEFAULT_TYPE):
    clear_history(update.effective_user.id)
    await update.message.reply_text("История очищена.", reply_markup=main_keyboard())

# ── Подписка / оплата через Telegram Stars ────────────────────────────────────

async def cmd_subscribe(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает текущий статус подписки + кнопки для покупки Basic / Pro."""
    uid = update.effective_user.id
    ensure_user(uid)
    plan = get_user_plan(uid)

    # Текущий статус
    if plan == PLAN_EXPIRED:
        status = "⏸ *Доступ закончился.* Оформи подписку чтобы продолжить."
    else:
        cfg = PLANS[plan]
        left = days_left(uid)
        status_lines = [f"*Твой тариф: {cfg['title']}*"]
        if left >= 0:
            status_lines.append(f"_Осталось дней: {left}_")
        status_lines.append("")
        status_lines.append(
            f"Сегодня использовано: {get_usage(uid,'msg')}/{cfg['msg_daily']} сообщений · "
            f"{get_usage(uid,'voice')}/{cfg['voice_daily']} голос · "
            f"{get_usage(uid,'photo')}/{cfg['photo_daily']} фото"
        )
        status = "\n".join(status_lines)

    lines = [status, "", "*Что можно выбрать:*"]
    for key in ("basic", "pro"):
        p = PLANS[key]
        lines.append(f"\n⭐ *{p['title']}* — {p['price_stars']} Stars / 30 дней")
        lines.append(p['description'])

    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton(f"Оформить Базовый · {PLANS['basic']['price_stars']}⭐",
                              callback_data="buy_basic")],
        [InlineKeyboardButton(f"Оформить Pro · {PLANS['pro']['price_stars']}⭐",
                              callback_data="buy_pro")],
    ])
    await update.message.reply_text("\n".join(lines), parse_mode="Markdown", reply_markup=kb)

async def send_stars_invoice(context: ContextTypes.DEFAULT_TYPE, chat_id: int, plan_key: str):
    """Отправляет инвойс Telegram Stars. provider_token для Stars должен быть пустой строкой."""
    cfg = PLANS.get(plan_key)
    if not cfg or cfg["price_stars"] <= 0:
        return
    from telegram import LabeledPrice
    await context.bot.send_invoice(
        chat_id=chat_id,
        title=f"Нова — {cfg['title']} (30 дней)",
        description=cfg["description"],
        payload=f"plan:{plan_key}",
        provider_token="",  # пусто = Telegram Stars
        currency="XTR",
        prices=[LabeledPrice(label=cfg["title"], amount=cfg["price_stars"])],
    )

async def handle_pre_checkout(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Подтверждаем платёж до списания. Telegram ждёт быстрый ответ."""
    q = update.pre_checkout_query
    try:
        payload = q.invoice_payload or ""
        if payload.startswith("plan:") and payload.split(":", 1)[1] in PLANS:
            await q.answer(ok=True)
        else:
            await q.answer(ok=False, error_message="Не удалось распознать тариф. Попробуй /subscribe заново.")
    except Exception as e:
        logging.error(f"pre_checkout error: {e}")
        try:
            await q.answer(ok=False, error_message="Техническая ошибка. Попробуй позже.")
        except Exception:
            pass

async def handle_successful_payment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Активирует подписку после успешной оплаты Stars."""
    msg = update.message
    sp = msg.successful_payment
    uid = update.effective_user.id
    try:
        payload = sp.invoice_payload or ""
        plan_key = payload.split(":", 1)[1] if payload.startswith("plan:") else None
        if plan_key not in PLANS:
            await msg.reply_text("Платёж прошёл, но тариф не распознался. Напиши в поддержку.")
            return
        activate_plan(uid, plan_key, stars_paid=sp.total_amount)
        cfg = PLANS[plan_key]
        await msg.reply_text(
            f"✨ Спасибо! Тариф *{cfg['title']}* активирован на 30 дней.\n"
            f"Лимиты: {cfg['msg_daily']} сообщений, {cfg['voice_daily']} голос, {cfg['photo_daily']} фото в день.",
            parse_mode="Markdown", reply_markup=main_keyboard())
    except Exception as e:
        logging.error(f"successful_payment error uid={uid}: {e}")
        await msg.reply_text("Платёж получен, но возникла ошибка активации. Напиши /subscribe — мы разберёмся.")

# ── Траты: команда /finance ───────────────────────────────────────────────────

async def cmd_finance(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    s = get_expenses_summary(uid, days=30)
    if s["count"] == 0:
        await update.message.reply_text(
            "💳 Пока нет записанных трат.\n\n"
            "Пришли мне фото чека — я сама извлеку сумму и категорию и добавлю в учёт.",
            reply_markup=main_keyboard())
        return
    lines = [f"💳 *Траты за 30 дней:* {s['total']:.0f} ₽  _({s['count']} операций)_", ""]
    for cat, amt in sorted(s["by_category"].items(), key=lambda x: -x[1]):
        lines.append(f"• {cat}: {amt:.0f} ₽")
    await send_safe(update, "\n".join(lines), main_keyboard())

# ── GDPR / приватность ────────────────────────────────────────────────────────

async def cmd_export(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Выгружает все данные юзера в текстовый файл и отправляет ему."""
    uid = update.effective_user.id
    lines = [f"# Экспорт данных пользователя {uid}", f"# Дата: {datetime.now().isoformat()}", ""]

    profile = get_profile(uid)
    if profile:
        lines.append("## Профиль")
        for k, v in profile.items():
            lines.append(f"- {k}: {v}")
        lines.append("")

    tasks = db_fetch("SELECT text, priority, timeframe, done, due_date, created_at FROM tasks WHERE user_id=?", (uid,))
    if tasks:
        lines.append(f"## Задачи ({len(tasks)})")
        for t in tasks:
            status = "✅" if t[3] else "⬜"
            lines.append(f"{status} [{t[1]}/{t[2]}] {t[0]}  _(создано: {t[5] or '?'}, срок: {t[4] or '—'})_")
        lines.append("")

    goals = db_fetch("SELECT text, progress, done, created_at FROM goals WHERE user_id=?", (uid,))
    if goals:
        lines.append(f"## Цели ({len(goals)})")
        for g in goals:
            lines.append(f"- {g[0]} — {g[1]}% {'(выполнено)' if g[2] else ''}")
        lines.append("")

    ideas = db_fetch("SELECT text, created_at FROM ideas WHERE user_id=?", (uid,))
    if ideas:
        lines.append(f"## Идеи ({len(ideas)})")
        for i in ideas:
            lines.append(f"- {i[0]}")
        lines.append("")

    msgs = db_fetch("SELECT role, content, created_at FROM messages WHERE user_id=? ORDER BY id", (uid,))
    if msgs:
        lines.append(f"## Переписка ({len(msgs)} сообщений)")
        for role, content, created in msgs:
            lines.append(f"\n### [{created}] {role}")
            lines.append(content)

    exp = db_fetch("SELECT amount, category, note, created_at FROM expenses WHERE user_id=?", (uid,))
    if exp:
        lines.append(f"\n## Траты ({len(exp)})")
        for amt, cat, note, created in exp:
            lines.append(f"- [{created}] {amt} ₽ — {cat} ({note or '—'})")

    data = "\n".join(lines).encode("utf-8")
    bio = io.BytesIO(data)
    bio.name = f"nova_export_{uid}.txt"
    await update.message.reply_document(document=bio, filename=bio.name,
                                        caption="📦 Здесь все твои данные в одном файле.")

async def cmd_delete_me(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Запрашивает подтверждение полного удаления данных."""
    uid = update.effective_user.id
    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("✅ Да, удалить всё",  callback_data="delete_me_yes"),
        InlineKeyboardButton("Отмена",              callback_data="delete_me_no"),
    ]])
    await update.message.reply_text(
        "⚠️ Ты действительно хочешь удалить ВСЕ свои данные?\n\n"
        "Будет удалено: профиль, задачи, цели, идеи, дневник, история переписки, "
        "привычки, траты, подписка, все настройки.\n\n"
        "*Это действие необратимо.*",
        parse_mode="Markdown", reply_markup=kb)

# ── /city — указать свой город (для часового пояса и будущих фич) ────────────

async def cmd_city(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Запрашивает у пользователя город и сохраняет в profile.city.
    Использование: /city или /city Москва (сразу с аргументом)."""
    uid = update.effective_user.id
    ensure_user(uid)
    args = context.args if hasattr(context, "args") else []
    if args:
        city = " ".join(args).strip()
        profile = get_profile(uid)
        profile["city"] = city
        save_profile(uid, profile)
        await update.message.reply_text(
            f"Запомнила: *{city}*. Теперь Нова знает где ты находишься.",
            parse_mode="Markdown", reply_markup=main_keyboard())
        return
    # Без аргумента — подсказка
    await update.message.reply_text(
        "В каком ты городе? Напиши: `/city Москва` (или твой город).\n\n"
        "Это нужно чтобы уведомления приходили в твоё местное время.",
        parse_mode="Markdown", reply_markup=main_keyboard())

# ── /forget — удаление конкретной памяти (Mem0 и история) ────────────────────

async def cmd_forget(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Забыть то, что Нова запомнила. Варианты: короткая память (диалог),
    долгая память (Mem0), или выборочно."""
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("🧹 Забыть только этот разговор",       callback_data="forget_chat")],
        [InlineKeyboardButton("🌫 Забыть всё что помнишь обо мне",     callback_data="forget_memory")],
        [InlineKeyboardButton("Отмена", callback_data="noop")],
    ])
    await update.message.reply_text(
        "Что мне забыть?\n\n"
        "• *Только этот разговор* — стерётся недавняя история сообщений. "
        "Профиль, задачи и цели останутся.\n"
        "• *Всё что помню о тебе* — уберу долгосрочную память (то что ты рассказывала ранее). "
        "Профиль, задачи, цели тоже сохранятся.",
        parse_mode="Markdown", reply_markup=kb)

# ── /backup — экспорт БД (только для владельца) ──────────────────────────────

async def cmd_backup(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Отправляет владельцу бота дамп локальной БД как .db файл.
    На Railway БД — временный SQLite; основной источник правды — Turso.
    Эта команда полезна для локальной проверки или быстрого снимка."""
    uid = update.effective_user.id
    if not OWNER_ID or uid != OWNER_ID:
        await update.message.reply_text("Эта команда только для владельца бота.")
        return
    import shutil
    try:
        src = "assistant.db"
        if not os.path.exists(src):
            await update.message.reply_text("Локальной assistant.db нет. Основные данные в Turso.")
            return
        snapshot = f"backup_{datetime.now().strftime('%Y%m%d_%H%M')}.db"
        shutil.copyfile(src, snapshot)
        with open(snapshot, "rb") as f:
            await update.message.reply_document(document=f, filename=snapshot,
                                                caption=f"💾 Бэкап БД на {datetime.now():%d.%m.%Y %H:%M}")
        try:
            os.remove(snapshot)
        except Exception:
            pass
    except Exception as e:
        logging.error(f"Backup failed: {e}")
        await update.message.reply_text(f"Не удалось сделать бэкап: {e}")

# ── Реферальная программа ────────────────────────────────────────────────────

def record_referral(invited_uid: int, inviter_uid: int):
    """Фиксирует кто кого пригласил. Один раз — если запись уже есть, не трогаем."""
    if invited_uid == inviter_uid:
        return
    existing = db_fetchone("SELECT inviter_user_id FROM referrals WHERE invited_user_id=?",
                           (invited_uid,))
    if existing:
        return
    db_exec("INSERT INTO referrals (invited_user_id, inviter_user_id, created_at) VALUES (?,?,?)",
            (invited_uid, inviter_uid, datetime.now().isoformat()))

async def cmd_invite(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает персональную реферальную ссылку.
    Кто перейдёт и зарегистрируется — попадёт в таблицу referrals.
    Когда монетизация включена, за приведённого друга можно выдавать бонусные дни."""
    uid = update.effective_user.id
    # Узнаём username бота из контекста
    try:
        me = await context.bot.get_me()
        bot_username = me.username
    except Exception:
        bot_username = "your_bot"
    link = f"https://t.me/{bot_username}?start=ref_{uid}"
    total = db_fetchone("SELECT COUNT(*) FROM referrals WHERE inviter_user_id=?", (uid,))
    count = total[0] if total else 0
    await update.message.reply_text(
        f"🔗 *Твоя пригласительная ссылка:*\n`{link}`\n\n"
        f"Пригласила уже: *{count}*\n\n"
        f"Отправь другу — если он зарегистрируется через твою ссылку, "
        f"я это запомню. Когда запустится оплата, за каждого друга получишь бонусные дни.",
        parse_mode="Markdown", reply_markup=main_keyboard())

# ── Промокоды ────────────────────────────────────────────────────────────────

def create_promo(code: str, plan: str = "basic", days: int = 30, max_uses: int = 1,
                 expires_at: str | None = None):
    """Создать промокод. Используется в /admin для владельца или вручную."""
    db_exec("""INSERT OR REPLACE INTO promo_codes (code, plan, days, max_uses, uses, expires_at, created_at)
               VALUES (?,?,?,?,0,?,?)""",
            (code.upper(), plan, days, max_uses, expires_at, datetime.now().isoformat()))

async def cmd_createpromo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """OWNER_ID может создать промокод прямо из чата.
    Формат: /createpromo КОД [plan=basic] [days=30] [uses=1]
    Пример: /createpromo ВЕСНА2026 basic 30 100  (код ВЕСНА2026, basic, 30 дней, до 100 применений)"""
    uid = update.effective_user.id
    if not OWNER_ID or uid != OWNER_ID:
        await update.message.reply_text("Команда только для владельца бота.")
        return
    args = context.args if hasattr(context, "args") else []
    if not args:
        await update.message.reply_text(
            "Формат: `/createpromo КОД [plan] [days] [max_uses]`\n"
            "Пример: `/createpromo ВЕСНА2026 basic 30 100`\n"
            "По умолчанию: plan=basic, days=30, max_uses=1",
            parse_mode="Markdown")
        return
    code = args[0].upper()
    plan = args[1] if len(args) > 1 else "basic"
    try:
        days = int(args[2]) if len(args) > 2 else 30
        max_uses = int(args[3]) if len(args) > 3 else 1
    except ValueError:
        await update.message.reply_text("days и max_uses должны быть числами.")
        return
    if plan not in PLANS:
        await update.message.reply_text(f"Неизвестный тариф «{plan}». Доступны: {', '.join(PLANS.keys())}")
        return
    create_promo(code, plan=plan, days=days, max_uses=max_uses)
    await update.message.reply_text(
        f"✅ Промокод *{code}* создан\n"
        f"Тариф: {PLANS[plan]['title']}\n"
        f"Срок: {days} дней\n"
        f"Макс. использований: {max_uses}",
        parse_mode="Markdown")

async def cmd_applycode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Применить промокод: /applycode ВЕСНА2026 — даст дни подписки."""
    uid = update.effective_user.id
    ensure_user(uid)
    args = context.args if hasattr(context, "args") else []
    if not args:
        await update.message.reply_text(
            "Чтобы применить промокод, напиши: `/applycode <код>`",
            parse_mode="Markdown", reply_markup=main_keyboard())
        return
    code = args[0].strip().upper()
    row = db_fetchone("SELECT plan, days, max_uses, uses, expires_at FROM promo_codes WHERE code=?", (code,))
    if not row:
        await update.message.reply_text("Такого промокода нет или он уже неактивен.",
                                        reply_markup=main_keyboard())
        return
    plan, days, max_uses, uses, expires_at = row
    if expires_at:
        try:
            if datetime.fromisoformat(expires_at) < datetime.now():
                await update.message.reply_text("Промокод истёк.", reply_markup=main_keyboard())
                return
        except Exception:
            pass
    if uses >= max_uses:
        await update.message.reply_text("Промокод уже использован максимальное число раз.",
                                        reply_markup=main_keyboard())
        return
    already = db_fetchone("SELECT 1 FROM promo_redemptions WHERE user_id=? AND code=?", (uid, code))
    if already:
        await update.message.reply_text("Ты уже применяла этот промокод.", reply_markup=main_keyboard())
        return
    # Даём дни: продлеваем подписку
    now = datetime.now()
    current = db_fetchone("SELECT valid_until FROM subscriptions WHERE user_id=?", (uid,))
    base = now
    if current and current[0]:
        try:
            dt = datetime.fromisoformat(current[0])
            if dt > now:
                base = dt
        except Exception:
            pass
    new_until = (base + timedelta(days=days)).isoformat()
    db_exec("""INSERT INTO subscriptions (user_id, plan, valid_until, last_payment_stars, last_payment_at)
               VALUES (?,?,?,0,?)
               ON CONFLICT(user_id) DO UPDATE SET plan=?, valid_until=?""",
            (uid, plan, new_until, now.isoformat(), plan, new_until))
    db_exec("INSERT INTO promo_redemptions (user_id, code, redeemed_at) VALUES (?,?,?)",
            (uid, code, now.isoformat()))
    db_exec("UPDATE promo_codes SET uses = uses + 1 WHERE code=?", (code,))
    await update.message.reply_text(
        f"✨ Промокод активирован! Тариф *{PLANS.get(plan, {}).get('title', plan)}* на *{days}* дней.",
        parse_mode="Markdown", reply_markup=main_keyboard())

# ── Генерация изображений (Pro-фича) ─────────────────────────────────────────

async def call_replicate_flux(prompt: str) -> bytes | None:
    """Генерирует картинку через Replicate FLUX.1 Schnell. Возвращает JPEG bytes или None."""
    if not REPLICATE_API_TOKEN:
        return None
    try:
        headers = {
            "Authorization": f"Bearer {REPLICATE_API_TOKEN}",
            "Content-Type": "application/json",
        }
        async with httpx.AsyncClient() as client:
            # Создаём prediction (FLUX.1 Schnell — быстрая модель, 4 шага, ~2 сек)
            r = await client.post(
                "https://api.replicate.com/v1/models/black-forest-labs/flux-schnell/predictions",
                headers=headers, timeout=60,
                json={"input": {"prompt": prompt, "num_outputs": 1,
                                 "aspect_ratio": "1:1", "output_format": "jpg"}})
            data = r.json()
            prediction_id = data.get("id")
            if not prediction_id:
                logging.error(f"Replicate: no prediction id: {data}")
                return None
            # Polling до готовности (обычно 2-5 сек)
            for _ in range(30):
                await asyncio.sleep(1)
                rr = await client.get(
                    f"https://api.replicate.com/v1/predictions/{prediction_id}",
                    headers=headers, timeout=15)
                pred = rr.json()
                status = pred.get("status")
                if status == "succeeded":
                    output = pred.get("output")
                    if isinstance(output, list) and output:
                        img_url = output[0]
                    elif isinstance(output, str):
                        img_url = output
                    else:
                        return None
                    # Скачиваем картинку
                    img_r = await client.get(img_url, timeout=30)
                    return img_r.content
                if status == "failed" or status == "canceled":
                    logging.error(f"Replicate failed: {pred.get('error')}")
                    return None
            return None
    except Exception as e:
        logging.error(f"Replicate error: {e}")
        return None

async def cmd_draw(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Сгенерировать картинку по тексту. /draw <описание>"""
    uid = update.effective_user.id
    args = context.args if hasattr(context, "args") else []
    if not args:
        await update.message.reply_text(
            "Опиши что нарисовать: `/draw закат над океаном, минимализм`",
            parse_mode="Markdown", reply_markup=main_keyboard())
        return
    if not REPLICATE_API_TOKEN:
        await update.message.reply_text(
            "Генерация картинок пока не настроена. Владельцу бота нужно добавить "
            "REPLICATE_API_TOKEN в настройки.",
            reply_markup=main_keyboard())
        return
    # Rate-limit отдельно — картинка стоит денег
    allowed, _ = check_rate_limit(uid)
    if not allowed:
        await update.message.reply_text("Подожди секунду после прошлого запроса 🙏")
        return
    prompt = " ".join(args).strip()
    await update.message.reply_text("🎨 Рисую... обычно 3-5 секунд")
    img_bytes = await call_replicate_flux(prompt)
    if not img_bytes:
        await update.message.reply_text("Не получилось сгенерировать( Попробуй другой запрос.",
                                        reply_markup=main_keyboard())
        return
    bio = io.BytesIO(img_bytes); bio.name = "nova_draw.jpg"
    await update.message.reply_photo(photo=bio, caption=f"🎨 _{prompt}_", parse_mode="Markdown",
                                     reply_markup=main_keyboard())

# ── Генерация PPTX-презентации ────────────────────────────────────────────────

async def cmd_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Быстрая генерация .pptx по теме. /presentation <тема>"""
    uid = update.effective_user.id
    args = context.args if hasattr(context, "args") else []
    if not args:
        await update.message.reply_text(
            "Напиши тему: `/presentation мои цели на 2026 год`",
            parse_mode="Markdown", reply_markup=main_keyboard())
        return
    topic = " ".join(args).strip()
    await update.message.reply_text("📊 Готовлю презентацию... 30-60 секунд")
    # 1) Просим Claude сгенерировать структуру слайдов (JSON)
    profile = get_profile(uid)
    system = build_system(profile, uid=uid)
    prompt = (
        f"Пользователь попросил презентацию по теме: «{topic}».\n\n"
        "Сгенерируй 5-7 слайдов как JSON-массив. Формат каждого слайда:\n"
        '{"title": "Заголовок", "bullets": ["пункт 1", "пункт 2", "пункт 3"]}\n\n'
        "Первый слайд — титул (без bullets, только title). Последний — вдохновляющее завершение.\n"
        "Верни ТОЛЬКО чистый JSON-массив, без комментариев и тегов, без markdown-обрамления."
    )
    try:
        raw = await call_claude([{"role": "user", "content": prompt}], system,
                                model=MODEL_SMART, max_tokens=MAX_TOKENS_REVIEW)
    except Exception as e:
        logging.error(f"Presentation claude error: {e}")
        await update.message.reply_text("Не получилось сгенерировать структуру(", reply_markup=main_keyboard())
        return
    # Парсим JSON
    m = re.search(r'\[.*\]', raw, re.DOTALL)
    if not m:
        await update.message.reply_text("Не смогла разобрать структуру слайдов. Попробуй другую формулировку.",
                                        reply_markup=main_keyboard())
        return
    try:
        slides = json.loads(m.group(0))
    except Exception as e:
        logging.error(f"Presentation JSON parse: {e}")
        await update.message.reply_text("Не удалось сформировать презентацию(", reply_markup=main_keyboard())
        return
    # 2) Собираем .pptx
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        await update.message.reply_text("Библиотека python-pptx не установлена.", reply_markup=main_keyboard())
        return
    prs = Presentation()
    # Титульный слайд
    if slides:
        first = slides[0]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = first.get("title", topic)
        if slide.placeholders and len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = "Нова · твой ассистент"
            except Exception:
                pass
    # Контентные слайды
    for sl in slides[1:]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sl.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        bullets = sl.get("bullets", [])
        if bullets:
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    bio.name = f"nova_presentation.pptx"
    await update.message.reply_document(document=bio, filename=bio.name,
                                        caption=f"📊 Презентация: _{topic}_",
                                        parse_mode="Markdown")

# ── NPS / отзыв (через 14 дней после регистрации) ────────────────────────────

async def check_feedback(context):
    """Раз в сутки (jobs) ищем юзеров, которым пора задать NPS-вопрос.
    Условие: прошло 14+ дней с первого сообщения, ещё не спрашивали."""
    cutoff = (datetime.now() - timedelta(days=14)).isoformat()
    # Юзеры у которых есть messages старше 14 дней и НЕТ записи в feedback
    candidates = db_fetch("""
        SELECT DISTINCT m.user_id
        FROM messages m
        LEFT JOIN feedback f ON f.user_id = m.user_id
        WHERE m.created_at <= ?
          AND f.user_id IS NULL
        LIMIT 50""", (cutoff,))
    for (uid,) in candidates:
        kb = InlineKeyboardMarkup([
            [InlineKeyboardButton(str(i), callback_data=f"nps_{i}") for i in range(1, 6)],
            [InlineKeyboardButton(str(i), callback_data=f"nps_{i}") for i in range(6, 11)],
        ])
        try:
            await context.bot.send_message(uid,
                "Мы с тобой уже две недели вместе 🌿\n\n"
                "На сколько оцениваешь меня от 1 до 10? Просто нажми цифру.\n"
                "_(1 — «удали, пожалуйста», 10 — «невозможно без Новы»)_",
                parse_mode="Markdown", reply_markup=kb)
            db_exec("INSERT OR IGNORE INTO feedback (user_id, asked_at) VALUES (?,?)",
                    (uid, datetime.now().isoformat()))
            await asyncio.sleep(0.3)
        except Exception as e:
            logging.warning(f"NPS send failed {uid}: {e}")

# ── Отложенные темы (парковка) ───────────────────────────────────────────────

async def cmd_parking(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает темы, которые Нова отложила во время онбординга.
    Пользователь выбирает одну — и Нова возвращается к разбору."""
    uid = update.effective_user.id
    topics = get_parking_topics(uid, only_open=True)
    if not topics:
        await update.message.reply_text(
            "Отложенных тем нет — всё актуальное обсудили 🌿",
            reply_markup=main_keyboard())
        return
    kb_rows = []
    for tid, topic, _ in topics[:8]:
        label = topic[:50]
        kb_rows.append([InlineKeyboardButton(f"💬 {label}", callback_data=f"discuss_{tid}")])
    kb_rows.append([InlineKeyboardButton("❌ Забыть всё это", callback_data="parking_clear")])
    await update.message.reply_text(
        "*Темы, которые мы отложили:*\nВыбери ту, о которой хочешь поговорить сейчас.",
        parse_mode="Markdown",
        reply_markup=InlineKeyboardMarkup(kb_rows))

# ── Админ / диагностика ──────────────────────────────────────────────────────

async def cmd_myid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Быстро узнать свой user_id — нужно при первой настройке OWNER_ID."""
    await update.message.reply_text(f"Твой user_id: `{update.effective_user.id}`", parse_mode="Markdown")

async def cmd_test_morning(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Проверить как выглядит утреннее приветствие, не дожидаясь 8:00.
    Доступно только OWNER_ID. НЕ трогает last_morning_sent — чтобы
    настоящая рассылка всё равно прошла утром."""
    uid = update.effective_user.id
    if not OWNER_ID or uid != OWNER_ID:
        await update.message.reply_text("Команда только для владельца бота.")
        return
    profile = get_profile(uid)
    local_now = user_now(profile)
    await update.message.reply_text("🧪 Генерирую утреннее приветствие (тестовая отправка)...")
    try:
        text = await _build_morning_text(uid, profile, local_now)
        await update.message.reply_text(text, parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Ошибка: {e}")

async def cmd_test_evening(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Тестовая отправка вечернего приветствия. Только OWNER_ID."""
    uid = update.effective_user.id
    if not OWNER_ID or uid != OWNER_ID:
        await update.message.reply_text("Команда только для владельца бота.")
        return
    profile = get_profile(uid)
    local_now = user_now(profile)
    await update.message.reply_text("🧪 Генерирую вечернее приветствие (тестовая отправка)...")
    try:
        text = await _build_evening_text(uid, profile, local_now)
        await update.message.reply_text(text, parse_mode="Markdown")
    except Exception as e:
        await update.message.reply_text(f"Ошибка: {e}")

async def cmd_admin(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Админская статистика — только для OWNER_ID."""
    uid = update.effective_user.id
    if not OWNER_ID or uid != OWNER_ID:
        await update.message.reply_text("Эта команда только для владельца бота.")
        return
    total_users = db_fetchone("SELECT COUNT(*) FROM users")[0]
    active_today = db_fetchone(
        "SELECT COUNT(DISTINCT user_id) FROM messages WHERE created_at >= ?",
        ((datetime.now() - timedelta(days=1)).isoformat(),))[0]
    active_7d = db_fetchone(
        "SELECT COUNT(DISTINCT user_id) FROM messages WHERE created_at >= ?",
        ((datetime.now() - timedelta(days=7)).isoformat(),))[0]
    paid = db_fetchone(
        "SELECT COUNT(*), COALESCE(SUM(last_payment_stars),0) FROM subscriptions WHERE plan!='free' AND valid_until > ?",
        (datetime.now().isoformat(),))
    msgs_24h = db_fetchone(
        "SELECT COUNT(*) FROM messages WHERE created_at >= ? AND role='assistant'",
        ((datetime.now() - timedelta(days=1)).isoformat(),))[0]
    refs = db_fetchone("SELECT COUNT(*) FROM referrals")[0]
    nps = db_fetch("SELECT rating FROM feedback WHERE rating IS NOT NULL")
    nps_line = ""
    if nps:
        avg = sum(r[0] for r in nps) / len(nps)
        nps_line = f"\nСредний NPS: *{avg:.1f}/10* ({len(nps)} ответов)"

    # График роста юзеров по дням за последние 14 дней
    chart_bio = None
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from collections import Counter
        rows = db_fetch("""SELECT DATE(created_at) FROM messages
                           WHERE created_at >= ? AND role='assistant'
                           ORDER BY created_at""",
                        ((datetime.now() - timedelta(days=14)).isoformat(),))
        counter = Counter([r[0] for r in rows if r[0]])
        if counter:
            days = sorted(counter.keys())
            values = [counter[d] for d in days]
            fig, ax = plt.subplots(figsize=(9, 4))
            ax.bar(days, values, color="#6C8EBF")
            ax.set_title("Ответы Новы за 14 дней")
            ax.set_ylabel("Сообщений")
            plt.xticks(rotation=45, ha="right", fontsize=8)
            plt.tight_layout()
            bio = io.BytesIO()
            plt.savefig(bio, format="png", dpi=100)
            plt.close(fig)
            bio.seek(0); bio.name = "admin_chart.png"
            chart_bio = bio
    except Exception as e:
        logging.warning(f"Admin chart failed: {e}")

    text = (
        "*📊 Статистика Новы*\n\n"
        f"Всего юзеров: *{total_users}*\n"
        f"Активных за сутки: *{active_today}*\n"
        f"Активных за 7 дней: *{active_7d}*\n\n"
        f"Платных подписок: *{paid[0] if paid else 0}*\n"
        f"Всего Stars получено: *{paid[1] if paid else 0}⭐*\n"
        f"Рефералов зафиксировано: *{refs}*{nps_line}\n\n"
        f"Ответов Новы за 24ч: *{msgs_24h}*"
    )
    if chart_bio:
        await update.message.reply_photo(photo=chart_bio, caption=text, parse_mode="Markdown")
    else:
        await update.message.reply_text(text, parse_mode="Markdown")

async def cmd_profile(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    if not profile:
        await update.message.reply_text("Профиль пока пустой.", reply_markup=main_keyboard()); return
    labels = {"name":"Имя","occupation":"Работа","goals":"Цели","pain":"Что мешает",
              "satisfied":"Что хорошо","day_rhythm":"Ритм дня","timezone":"Часовой пояс","character":"Характер","notes":"Заметки"}
    lines = ["*Что я знаю о тебе:*\n"]
    for k, l in labels.items():
        if profile.get(k): lines.append(f"*{l}:* {profile[k]}")
    await send_safe(update, "\n".join(lines), main_keyboard())

async def cmd_today(update: Update, context: ContextTypes.DEFAULT_TYPE):
    tasks = get_today_tasks(update.effective_user.id)
    await send_safe(update, f"📅 *На сегодня:*\n\n{format_tasks(tasks)}", main_keyboard())

async def cmd_week(update: Update, context: ContextTypes.DEFAULT_TYPE):
    tasks = get_tasks(update.effective_user.id, timeframe="week")
    await send_safe(update, f"📆 *На неделю:*\n\n{format_tasks(tasks)}", main_keyboard())

async def cmd_plan(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    await send_safe(update, format_week_plan(uid), main_keyboard())

async def cmd_goals(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    goals = get_goals(uid)
    if goals:
        await send_safe(update, f"*🎯 Твои цели:*\n\n{format_goals(goals)}", main_keyboard())
    else:
        await update.message.reply_text("Целей пока нет 🎯", reply_markup=main_keyboard())

async def cmd_ideas(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ideas = get_ideas(uid)
    if ideas:
        lines = ["*💡 Идеи и мечты:*\n"] + [f"• [{i[0]}] {i[1]}" for i in ideas]
        await send_safe(update, "\n".join(lines), main_keyboard())
    else:
        await update.message.reply_text("Идей пока нет 💡", reply_markup=main_keyboard())

async def cmd_focus(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    tasks = get_tasks(uid)
    profile = get_profile(uid)
    system = build_system(profile, uid=uid)
    task_list = "\n".join([f"- ({t[2]}) {t[1]}" for t in tasks[:10]]) if tasks else "Задач нет"
    try:
        response = await call_claude(
            [{"role": "user", "content": f"Режим фокуса. Задачи:\n{task_list}\n\nОдна самая важная прямо сейчас — какая и почему?"}],
            system, model=MODEL_SMART)
        clean = await process_response(uid, response)
        await send_safe(update, clean, main_keyboard())
    except:
        await update.message.reply_text("Что-то пошло не так)", reply_markup=main_keyboard())

async def cmd_checkin(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    system = build_system(profile, uid=uid)
    try:
        response = await call_claude(
            [{"role": "user", "content": "Проведи короткий чекин — спроси как я себя чувствую и какая энергия."}],
            system, model=MODEL_SMART)
        clean = await process_response(uid, response)
        await send_safe(update, clean, main_keyboard())
    except:
        await update.message.reply_text("Как ты сейчас? 🙂", reply_markup=main_keyboard())

async def cmd_help(update: Update, context: ContextTypes.DEFAULT_TYPE):
    subscribe_block = "/subscribe — мой тариф и подписка\n" if PAYMENTS_ENABLED else ""
    text = f"""*Нова — твой личный ассистент* 🤖

*📋 Задачи и планирование*
/today — задачи на сегодня
/week — план на неделю
/month — задачи на месяц
/focus — самая важная задача прямо сейчас
/brain — выгрузи всё из головы, я разберу

*🎯 Цели и развитие*
/goals — цели и прогресс в %
/wins — твои победы и достижения
/review — ежемесячный разбор прогресса

*💡 Идеи и рефлексия*
/ideas — идеи и желания
/reflect — вопрос для самоанализа
/ask — честный коучинговый ответ
/journal — личный дневник

*🌀 Трекеры*
/mood — настроение
/energy — уровень энергии
/habits — трекер привычек
/sphere — колесо жизни по сферам

*🔧 Настройки*
/checkin — чекин состояния
/profile — мой профиль
/city — указать свой город
/settings — уведомления, часовой пояс, стиль
/calendar — подключить Google Календарь
/calshow — показать события в Calendar
/report — отчёт, графики, PDF

*🎨 Креатив*
/draw — сгенерировать картинку по описанию
/presentation — сделать .pptx по теме

*💳 Финансы*
/finance — траты за месяц (фото чека → учёт)
{subscribe_block}
*🤝 Друзьям*
/invite — моя ссылка-приглашение
/applycode — применить промокод

*🔒 Приватность*
/forget — забыть разговор или память обо мне
/parking — вернуться к отложенным темам
/export — скачать все мои данные
/delete\\_me — удалить меня и все данные

Пишу в любом формате — текст, голос, фото 🎤📸"""
    await send_safe(update, text, main_keyboard())

async def cmd_month(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    tasks = get_tasks(uid, timeframe="month")
    await send_safe(update, f"🗓 *На месяц:*\n\n{format_tasks(tasks)}", main_keyboard())

async def cmd_sphere(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    spheres_score = profile.get("spheres_score", "")
    if spheres_score:
        lines = ["*🌀 Колесо жизни:*\n"]
        for pair in spheres_score.split(","):
            if ":" in pair:
                k, v = pair.strip().split(":", 1)
                try:
                    val = max(0, min(10, int(v.strip())))
                except ValueError:
                    val = 5
                bar = "█" * val + "░" * (10 - val)
                lines.append(f"{k.strip().capitalize()}: {bar} {val}/10")
        chart = generate_wheel_chart(uid)
        if chart:
            await context.bot.send_photo(uid, photo=chart)
        await send_safe(update, "\n".join(lines) + "\n\nОбновить оценки? Просто напиши мне новые.", main_keyboard())
    else:
        await send_safe(update, "Оценки по сферам ещё не заполнены. Пройди /start чтобы заполнить колесо жизни, или напиши мне оценки в свободной форме.", main_keyboard())

async def cmd_reflect(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    q = random.choice(REFLECT_QUESTIONS)
    profile = get_profile(uid)
    address = profile.get("address") or profile.get("name") or ""
    prefix = f"{address}, " if address else ""
    await send_safe(update, f"🪞 *Вопрос для размышления:*\n\n_{prefix}{q}_", main_keyboard())
    set_followup(uid)

async def cmd_wins(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    wins = get_wins(uid)
    if wins:
        lines = ["*🏆 Твои победы:*\n"]
        for w in wins:
            dt = w[1][:10] if w[1] else ""
            lines.append(f"• {w[0]}" + (f" _{dt}_" if dt else ""))
        await send_safe(update, "\n".join(lines), main_keyboard())
    else:
        await send_safe(update, "Побед пока нет — но это ненадолго 💪\nНапиши мне о любом своём достижении, я сохраню.", main_keyboard())

async def cmd_mood(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Как твоё настроение сейчас? Выбери от 1 до 10:",
        reply_markup=score_keyboard("mood"))

async def cmd_habits(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    habits = get_habits(uid)
    if not habits:
        await send_safe(update, "Привычек пока нет. Напиши мне какую привычку хочешь отслеживать — я добавлю.", main_keyboard())
        return
    await update.message.reply_text("*📌 Трекер привычек — сегодня:*",
                                     parse_mode="Markdown", reply_markup=habits_keyboard(uid))

async def cmd_energy(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "Уровень энергии сейчас — от 1 до 10:",
        reply_markup=score_keyboard("energy"))

async def cmd_journal(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    system = build_system(profile, uid=uid)
    entries = get_journal_entries(uid, limit=3)
    context_block = ""
    if entries:
        context_block = "\nПоследние записи:\n" + "\n".join([f"— {e[0]}: {e[1][:60]}…" for e in entries])
    try:
        response = await call_claude(
            [{"role": "user", "content": f"Задай один глубокий вопрос для дневниковой записи. Учитывай профиль.{context_block}"}],
            system, model=MODEL_SMART)
        context.user_data["journal_question"] = response
        await send_safe(update, f"📔 *Дневник*\n\n{response}", None)
    except:
        q = "Что сегодня было для тебя самым значимым?"
        context.user_data["journal_question"] = q
        await send_safe(update, f"📔 *Дневник*\n\n{q}", None)

async def cmd_brain(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await send_safe(update,
        "🧠 *Выгрузка мыслей*\n\nПиши всё подряд — дела, идеи, тревоги, планы, случайные мысли. Не фильтруй. Я сама разберу по категориям.",
        None)
    context.user_data["mode"] = "brain"

async def cmd_ask(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await send_safe(update,
        "🔍 *Честный анализ*\n\nЗадай мне вопрос о себе — и я отвечу честно, как зеркало. Без лишней мягкости.",
        None)
    context.user_data["mode"] = "ask"

async def cmd_review(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    tasks = get_tasks(uid)
    goals = get_goals(uid)
    mood_hist = get_mood_history(uid, days=30)
    energy_hist = get_energy_history(uid, days=30)
    stats = get_sphere_stats(uid)

    avg_mood = round(sum(r[0] for r in mood_hist) / len(mood_hist), 1) if mood_hist else "нет данных"
    avg_energy = round(sum(r[0] for r in energy_hist) / len(energy_hist), 1) if energy_hist else "нет данных"
    goals_block = "\n".join([f"• {g[1]} — {g[4]}%" for g in goals[:5]]) if goals else "нет"
    stats_block = ", ".join([f"{k}: {v}д" for k, v in list(stats.items())[:5]]) if stats else "нет данных"

    system = build_system(profile, uid=uid)
    prompt = f"""Сделай ежемесячный разбор для пользователя.

Открытых задач: {len(tasks)}
Активных целей:\n{goals_block}
Ср. настроение за месяц: {avg_mood}/10
Ср. энергия за месяц: {avg_energy}/10
Активность по сферам: {stats_block}

Структура:
1. Что работает — честно и конкретно
2. Что не работает — без осуждения
3. Паттерны которые ты замечаешь
4. Одна главная рекомендация на следующий месяц
5. Короткое вдохновляющее завершение

Стиль: глубокий, честный, поддерживающий."""

    await update.message.reply_text("Готовлю ежемесячный разбор... 📊")
    try:
        response = await call_claude([{"role": "user", "content": prompt}], system, model=MODEL_SMART)
        chart = generate_sphere_chart(uid)
        if chart:
            await context.bot.send_photo(uid, photo=chart, caption="📊 Активность и прогресс")
        await send_safe(update, response, main_keyboard())
    except Exception as e:
        logging.error(f"Review error: {e}")
        await update.message.reply_text("Что-то пошло не так(", reply_markup=main_keyboard())

async def cmd_settings(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    profile = get_profile(uid)
    tz = profile.get("timezone", "не задан")
    style = "подробно" if profile.get("info_style", "detailed") == "detailed" else "кратко"
    feedback = "мягкий" if profile.get("feedback_style", "soft") == "soft" else "прямой"
    text = f"*⚙️ Настройки*\n\nЧасовой пояс: {tz}\nСтиль ответов: {style}\nТон: {feedback}"
    await update.message.reply_text(text, parse_mode="Markdown", reply_markup=settings_keyboard(profile))

async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)
    allowed, reason = check_rate_limit(uid)
    if not allowed:
        await update.message.reply_text("Секунду, я ещё с прошлым сообщением 🙏")
        return
    # Проверка тарифа: голос — отдельная квота
    user_row = get_user(uid)
    onboarding_done = bool(user_row and user_row[1])
    if onboarding_done:
        ok, msg = check_plan_limit(uid, kind="voice")
        if not ok:
            await update.message.reply_text(msg)
            return
    clear_followup(uid)
    await update.message.reply_text("Слушаю... 🎤")
    try:
        if not os.environ.get("GROQ_API_KEY"):
            await update.message.reply_text(
                "⚠️ Голосовые сообщения не настроены.\n\nНужно добавить GROQ_API_KEY в Railway → Variables.\nПолучи бесплатный ключ на console.groq.com")
            return
        voice = update.message.voice
        file = await context.bot.get_file(voice.file_id)
        audio_bytes = await file.download_as_bytearray()
        text = await call_groq_voice(bytes(audio_bytes))
        if not text:
            await update.message.reply_text("Не смогла расшифровать( Попробуй ещё раз.")
            return
        profile = get_profile(uid)
        system = build_system(profile, onboarding_mode=not onboarding_done, uid=uid)
        history = get_history(uid)
        history.append({"role": "user", "content": text})
        save_msg(uid, "user", f"[голосовое] {text}")
        response = await call_claude(history, system)
        clean = await process_response(uid, response, skip_calendar=not onboarding_done)
        save_msg(uid, "assistant", clean)
        # Follow-up только в онбординге — см. комментарий в handle_message
        if not onboarding_done and "?" in clean:
            set_followup(uid)
        else:
            clear_followup(uid)
        if onboarding_done:
            bump_usage(uid, "voice")
        await send_safe(update, f"_Ты сказала:_ {text}\n\n{clean}", main_keyboard() if onboarding_done else onboarding_keyboard())
    except Exception as e:
        logging.error(f"Voice error uid={uid}: {e}")
        await update.message.reply_text("Не смогла обработать голосовое( Попробуй ещё раз.")

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)
    allowed, _ = check_rate_limit(uid)
    if not allowed:
        await update.message.reply_text("Секунду, я ещё с прошлым сообщением 🙏")
        return
    user = get_user(uid)
    onboarding_done = bool(user and user[1])
    if onboarding_done:
        ok, msg = check_plan_limit(uid, kind="photo")
        if not ok:
            await update.message.reply_text(msg)
            return
    clear_followup(uid)
    await update.message.reply_text("Смотрю фото... 👀")
    photo = update.message.photo[-1]
    file = await context.bot.get_file(photo.file_id)
    photo_bytes = await file.download_as_bytearray()
    image_b64 = base64.b64encode(bytes(photo_bytes)).decode('utf-8')
    caption = update.message.caption or ""
    profile = get_profile(uid)
    system = build_system(profile, onboarding_mode=not onboarding_done, uid=uid)

    # Подсказка модели: если на фото чек — нужно вернуть структурированный тег.
    # Формат тега: [EXPENSE: сумма | категория | короткое описание]
    # Примеры категорий: еда, транспорт, кафе, здоровье, дом, развлечения, одежда, прочее.
    prompt = (
        f"Пользователь прислал фото. {'Подпись: ' + caption if caption else ''} "
        "Опиши что видишь, извлеки задачи, планы, важную информацию. "
        "Если на фото ЧЕК или квитанция — обязательно добавь в конце отдельной строкой тег:\n"
        "[EXPENSE: <итоговая сумма числом> | <категория одним словом: еда|транспорт|кафе|здоровье|"
        "дом|развлечения|одежда|услуги|прочее> | <где потрачено / на что]\n"
        "Не выдумывай сумму — если её не видно, просто напиши [EXPENSE: 0 | прочее | не видно суммы]."
    )
    try:
        response = await call_claude_vision(image_b64, system, prompt)

        # Парсим тег трат, сохраняем в expenses и удаляем из ответа юзеру
        exp_added = None
        m = re.search(r"\[EXPENSE:\s*([0-9]+(?:[.,][0-9]+)?)\s*\|\s*([^|\]]+?)\s*\|\s*([^\]]+?)\]",
                      response, flags=re.IGNORECASE)
        clean_text = response
        if m:
            try:
                amount = float(m.group(1).replace(",", "."))
                category = m.group(2).strip().lower()
                note = m.group(3).strip()
                if amount > 0:
                    add_expense(uid, amount, category, note)
                    exp_added = (amount, category, note)
            except Exception as e:
                logging.warning(f"Expense parse failed: {e}")
            clean_text = re.sub(r"\[EXPENSE:[^\]]+\]", "", response).strip()

        clean = await process_response(uid, clean_text, skip_calendar=not onboarding_done)
        if exp_added:
            amt, cat, note = exp_added
            clean += f"\n\n💳 Записала трату: *{amt:g} ₽* — {cat} ({note})"
        save_msg(uid, "user", f"[фото] {caption}")
        save_msg(uid, "assistant", clean)
        if onboarding_done:
            bump_usage(uid, "photo")
        await send_safe(update, clean, main_keyboard() if onboarding_done else onboarding_keyboard())
    except Exception as e:
        logging.error(f"Photo error: {e}")
        await update.message.reply_text("Не смогла обработать фото( Попробуй ещё раз.")

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)
    allowed, _ = check_rate_limit(uid)
    if not allowed:
        await update.message.reply_text("Секунду, я ещё с прошлым сообщением 🙏")
        return
    clear_followup(uid)
    doc = update.message.document
    if not doc.mime_type or not doc.mime_type.startswith('text'):
        await update.message.reply_text("Пока умею читать только текстовые файлы (.txt, .md)")
        return
    await update.message.reply_text("Читаю документ... 📄")
    file = await context.bot.get_file(doc.file_id)
    file_bytes = await file.download_as_bytearray()
    try:
        text_content = file_bytes.decode('utf-8')
    except:
        text_content = file_bytes.decode('latin-1')
    if len(text_content) > 3000:
        text_content = text_content[:3000] + "...[обрезано]"
    profile = get_profile(uid)
    user = get_user(uid)
    system = build_system(profile, onboarding_mode=not user[1], uid=uid)
    history = get_history(uid)
    history.append({"role": "user", "content": f"Я прислала документ '{doc.file_name}':\n\n{text_content}\n\nПроанализируй, извлеки задачи и важную информацию."})
    save_msg(uid, "user", f"[документ: {doc.file_name}]")
    try:
        response = await call_claude(history, system)
        clean = await process_response(uid, response, skip_calendar=not user[1])
        save_msg(uid, "assistant", clean)
        await send_safe(update, clean, main_keyboard() if user[1] else onboarding_keyboard())
    except Exception as e:
        logging.error(f"Doc error: {e}")
        await update.message.reply_text("Не смогла обработать документ(")

async def handle_forward(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)
    clear_followup(uid)
    text = update.message.text or update.message.caption or ""
    if not text:
        await update.message.reply_text("Пересланное сообщение без текста — не могу обработать(")
        return
    profile = get_profile(uid)
    user = get_user(uid)
    system = build_system(profile, onboarding_mode=not user[1], uid=uid)
    history = get_history(uid)
    history.append({"role": "user", "content": f"Я переслала сообщение:\n\n{text}\n\nОбработай — извлеки задачи, важную информацию или просто прокомментируй."})
    save_msg(uid, "user", f"[пересланное] {text[:100]}")
    try:
        response = await call_claude(history, system)
        clean = await process_response(uid, response, skip_calendar=not user[1])
        save_msg(uid, "assistant", clean)
        await send_safe(update, clean, main_keyboard() if user[1] else onboarding_keyboard())
    except:
        await update.message.reply_text("Что-то пошло не так)")

async def handle_menu_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    text = update.message.text
    if text == "📋 Задачи":
        await update.message.reply_text("Выбери:", reply_markup=tasks_keyboard()); return True
    if text == "🎯 Цели":
        await update.message.reply_text("Выбери:", reply_markup=goals_keyboard()); return True
    if text == "🌀 Сферы жизни":
        await update.message.reply_text("Выбери сферу:", reply_markup=spheres_keyboard()); return True
    if text == "💡 Идеи":
        ideas = get_ideas(uid)
        if ideas:
            lines = ["*💡 Идеи и мечты:*\n"] + [f"• [{i[0]}] {i[1]}" for i in ideas]
            await send_safe(update, "\n".join(lines), main_keyboard())
        else:
            await update.message.reply_text("Идей пока нет... поделись)", reply_markup=main_keyboard())
        return True
    if text == "📊 Дашборд":
        await update.message.reply_text(format_dashboard(uid), reply_markup=main_keyboard()); return True
    if text == "📅 План недели":
        await send_safe(update, format_week_plan(uid), main_keyboard()); return True
    return False

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    ensure_user(uid)

    # Rate-limit: защита от цикла/спама — ДО любых LLM-вызовов.
    allowed, reason = check_rate_limit(uid)
    if not allowed:
        if reason == "too_fast":
            await update.message.reply_text("Секунду, я ещё отвечаю на прошлое... 🙏")
        else:
            await update.message.reply_text("Слишком много сообщений подряд. Давай сделаем паузу на минуту 🌿")
        return

    if getattr(update.message, 'forward_origin', None) or getattr(update.message, 'forward_from', None) or getattr(update.message, 'forward_from_chat', None):
        await handle_forward(update, context)
        return

    user = get_user(uid)
    if not user:
        await update.message.reply_text("Привет! Напиши /start чтобы начать.")
        return
    if await handle_menu_button(update, context):
        return

    text = update.message.text
    if not text:
        return

    profile = get_profile(uid)
    onboarding_done = user[1]

    # Если подписка/триал истекли — блокируем ВСЕГДА (даже если онбординг не завершён).
    # Иначе проверяем лимит только после онбординга — чтобы новый юзер мог пройти старт.
    plan_now = get_user_plan(uid)
    if plan_now == PLAN_EXPIRED:
        ok, msg = check_plan_limit(uid, kind="msg")
        await send_safe(update, msg, main_keyboard())
        return
    if onboarding_done:
        ok, msg = check_plan_limit(uid, kind="msg")
        if not ok:
            await send_safe(update, msg, main_keyboard())
            return

    # Загружаем события и список календарей если пользователь управляет Calendar
    cal_events = None
    _cal_keywords = ("удали", "очисти", "убери", "перенеси", "измени", "редактир",
                     "событи", "календар", "папк", "группу", "список")
    if get_google_token(uid) and any(kw in text.lower() for kw in _cal_keywords):
        try:
            cal_events = await list_calendar_events(uid, max_results=30, include_past=True)
            cals = await list_calendars(uid)
            if cals:
                cal_events = cal_events or []
                cal_events.append({"_calendars": cals})
        except Exception as e:
            logging.warning(f"Calendar context load error: {e}")

    system = build_system(profile, onboarding_mode=not onboarding_done, uid=uid, cal_events=cal_events)

    if onboarding_done:
        tasks = get_tasks(uid)
        if tasks:
            system += "\n\nАктуальные задачи:\n" + "\n".join([f"[{t[0]}] ({t[2]}) {t[1]}" for t in tasks[:10]])
        frozen = get_frozen_items(uid)
        if frozen and len(tasks) == 0:
            items_text = "\n".join([f"- {f[1]}" for f in frozen])
            system += f"\n\nЗамороженные идеи/цели (давно без движения):\n{items_text}\nЕсли уместно — предложи запланировать одну из них."

    # Подтверждение плана Google Calendar (CAL_PLAN)
    cal_pending = context.user_data.get("cal_pending")
    if cal_pending:
        text_lower = text.lower()
        confirm_words = ("да", "вноси", "всё верно", "все верно", "ок", "ok", "добавь",
                         "подходит", "согласна", "согласен", "давай", "отлично", "супер", "👍")
        reject_words = ("нет", "не надо", "отмена", "изменить", "подожди", "стоп", "нет спасибо")
        if any(w in text_lower for w in confirm_words):
            context.user_data.pop("cal_pending")
            added = []
            for parts in cal_pending:
                title   = parts[0] if len(parts) > 0 else ""
                date    = parts[1] if len(parts) > 1 else None
                t_time  = parts[2] if len(parts) > 2 else None
                prio    = parts[3] if len(parts) > 3 else "normal"
                notes   = parts[4] if len(parts) > 4 else ""
                if not title: continue
                add_task(uid, title, prio, "general", "week")
                date_str = await add_to_calendar(uid, title,
                    due_date=date or None, event_time=t_time or None,
                    priority=prio, description=notes)
                icon = "🔴" if prio == "urgent" else "🟡" if prio == "important" else "🔵"
                added.append(f"{icon} {title}" + (f" — {date_str}" if date_str else ""))
            msg = "✅ Готово! Добавила в Google Calendar:\n\n" + "\n".join(added) if added else "✅ Добавила в Google Calendar!"
            await send_safe(update, msg, main_keyboard())
            return
        elif any(w in text_lower for w in reject_words):
            context.user_data.pop("cal_pending")
            # Не возвращаемся — даём обычному разговору продолжиться

    # Секретные команды сброса
    if text.strip().lower() in ("полный сброс", "full reset"):
        wipe_user_data(uid)
        await update.message.reply_text("Полный сброс выполнен. Напиши /start")
        return
    if text.strip().lower() == "сброс истории":
        clear_history(uid)
        await update.message.reply_text("История диалога очищена.", reply_markup=main_keyboard())
        return

    # Режим выгрузки мыслей
    mode = context.user_data.get("mode")
    if mode == "brain":
        context.user_data.pop("mode", None)
        profile = get_profile(uid)
        system = build_system(profile, uid=uid)
        try:
            response = await call_claude(
                [{"role": "user", "content": f"Пользователь выгружает всё из головы. Разбери по категориям: задачи, идеи, тревоги, цели. Для каждой категории дай короткий список. Текст:\n\n{text}"}],
                system, model=MODEL_SMART)
            clean = await process_response(uid, response)
            save_msg(uid, "user", f"[brain dump] {text[:100]}")
            save_msg(uid, "assistant", clean)
            await send_safe(update, clean, main_keyboard())
        except:
            await update.message.reply_text("Что-то пошло не так(", reply_markup=main_keyboard())
        return
    # Режим честного анализа
    if mode == "ask":
        context.user_data.pop("mode", None)
        profile = get_profile(uid)
        system = build_system(profile, uid=uid)
        try:
            response = await call_claude(
                [{"role": "user", "content": f"Пользователь задаёт вопрос о себе и просит честный коучинговый ответ — как зеркало, без лишней мягкости, но с уважением. Вопрос: {text}"}],
                system, model=MODEL_SMART)
            clean = await process_response(uid, response)
            save_msg(uid, "user", f"[ask] {text}")
            save_msg(uid, "assistant", clean)
            await send_safe(update, clean, main_keyboard())
        except:
            await update.message.reply_text("Что-то пошло не так(", reply_markup=main_keyboard())
        return
    # Режим дневника — ответ на вопрос
    if "journal_question" in context.user_data:
        question = context.user_data.pop("journal_question")
        save_journal(uid, question, text)
        await send_safe(update, "Записала в дневник 📔", main_keyboard())
        return
    # Добавление привычки
    if mode == "habit_add":
        context.user_data.pop("mode", None)
        get_or_create_habit(uid, text.strip())
        await send_safe(update, f"Привычка «{text.strip()}» добавлена! Отмечай каждый день через /habits 💪", main_keyboard())
        return
    # Установка часового пояса через текст
    if mode == "set_tz":
        context.user_data.pop("mode", None)
        profile = get_profile(uid)
        system = build_system(profile, uid=uid)
        try:
            response = await call_claude(
                [{"role": "user", "content": f"Пользователь написал своё текущее время: '{text}'. Вычисли UTC offset и ответь только строкой вида UTC+N или UTC-N."}],
                system)
            tz_str = response.strip().split()[0]
            profile["timezone"] = tz_str
            save_profile(uid, profile)
            await send_safe(update, f"Часовой пояс обновлён: {tz_str}", main_keyboard())
        except:
            await update.message.reply_text("Не смог вычислить часовой пояс. Напиши в формате UTC+3", reply_markup=main_keyboard())
        return

    # Детектим когда пользователь говорит о победе/достижении
    win_keywords = ("сделал", "завершил", "закончил", "выполнил", "достиг", "получил", "удалось", "победил", "справился")
    if any(kw in text.lower() for kw in win_keywords) and len(text) < 200:
        context.user_data["potential_win"] = text

    clear_followup(uid)

    history = get_history(uid)
    history.append({"role": "user", "content": text})
    save_msg(uid, "user", text)

    # Обогащаем системный промпт релевантными воспоминаниями из Mem0
    mem_block = await mem0_search(uid, text)
    enriched_system = system + (f"\n\n{mem_block}" if mem_block else "")

    try:
        # Во время онбординга — умная модель (Sonnet) + больший запас токенов
        # на развёрнутые разборы. В обычном диалоге — авто-выбор (часто DeepSeek).
        if not onboarding_done:
            response = await call_claude(history, enriched_system,
                                         model=MODEL_SMART, max_tokens=MAX_TOKENS_ONBOARD)
        else:
            response = await call_claude(history, enriched_system,
                                         max_tokens=MAX_TOKENS_DEFAULT)
    except Exception as e:
        logging.error(f"Error: {e}")
        await update.message.reply_text("Что-то пошло не так... попробуй ещё раз)"); return

    cal_plan_buffer = []
    clean = await process_response(uid, response, skip_calendar=not onboarding_done, cal_plan_buffer=cal_plan_buffer)

    # Если Claude предложил план для календаря — сохраняем для подтверждения
    if cal_plan_buffer:
        context.user_data["cal_pending"] = cal_plan_buffer

    save_msg(uid, "assistant", clean)
    # Follow-up ставим ТОЛЬКО во время онбординга — чтобы помочь человеку
    # закончить первичную настройку, если он ушёл из чата посередине.
    # В обычном диалоге НЕ переспрашиваем: человек сам вернётся, если захочет.
    if not onboarding_done and "?" in clean:
        set_followup(uid)
    else:
        # На случай, если followup остался с прошлого диалога — чистим,
        # чтобы бот не задавал вопрос вдогонку уже завершённой теме.
        clear_followup(uid)

    # Сохраняем диалог в Mem0 для долгосрочной памяти
    await mem0_add(uid, [
        {"role": "user", "content": text},
        {"role": "assistant", "content": clean},
    ])

    # Предлагаем отметить победу если Нова закрыла задачу
    if "potential_win" in context.user_data:
        if any(kw in clean.lower() for kw in ("выполнен", "закрыт", "готово", "сделан", "✅")):
            win_text = context.user_data.pop("potential_win")
            add_win(uid, win_text)
        else:
            context.user_data.pop("potential_win", None)

    # Счётчик лимита инкрементируем ПОСЛЕ успешной обработки — чтобы сбой
    # в Claude не крал попытки.
    if onboarding_done:
        bump_usage(uid, "msg")

    await send_safe(update, clean, main_keyboard() if onboarding_done else onboarding_keyboard())

# ── Вспомогательные: защита от дублирующих рассылок ──────────────────────────
# Причина проблемы с «тремя добрыми вечерами»: каждый перезапуск контейнера
# на Railway создавал job_queue заново, и если местное время пользователя
# всё ещё попадало в окно (8:00 для утра, 19:00 для вечера), сообщение
# отправлялось ПОВТОРНО. Теперь сохраняем в профиле дату последней отправки
# и пропускаем, если уже отправили сегодня.

def _mark_sent(uid: int, field: str, day_iso: str):
    """Сохранить дату отправки в profile.<field> (атомарно через save_profile)."""
    profile = get_profile(uid)
    profile[field] = day_iso
    save_profile(uid, profile)

def _already_sent_today(profile: dict, field: str, day_iso: str) -> bool:
    return profile.get(field) == day_iso

async def _build_morning_text(uid: int, profile: dict, local_now: datetime) -> str:
    """Генерирует текст утреннего сообщения (без отправки).
    Выделено в отдельную функцию, чтобы команда /test_morning могла
    переиспользовать ту же логику."""
    address = profile.get("address") or profile.get("name") or ""
    today_tasks = get_today_tasks(uid)
    urgent = get_tasks(uid, priority="urgent")
    goals = get_goals(uid)
    quote_text, quote_author = get_random_quote(uid)
    notif_extras = profile.get("notif_extras", "")

    task_block = ""
    if today_tasks:
        task_block = f"Задачи на сегодня:\n" + "\n".join([f"• {t[1]}" for t in today_tasks[:5]])
    elif urgent:
        task_block = f"Срочных задач: {len(urgent)}"

    goal_block = ""
    if goals:
        goal_block = "Активных целей: " + str(len(goals))

    system = build_system(profile, uid=uid)
    prompt = f"""Сгенерируй утреннее уведомление для пользователя. Используй эти данные:

Обращение: {address}
Дата: {local_now.strftime('%d.%m.%Y, %A')}
Цитата дня: «{quote_text}» — {quote_author}
{task_block}
{goal_block}
{f"Дополнительно включи: {notif_extras}" if notif_extras else ""}

Структура (одним цельным сообщением, без лишних блоков):
1. Тёплое короткое приветствие с датой
2. Цитата дня — курсивом, с автором
3. Напоминание о задачах на сегодня (если есть) — коротко, по делу
4. Одна заряжающая фраза-мотивация на день

ВАЖНО:
- НЕ задавай никаких вопросов пользователю — это приветствие, а не диалог
- Не вставляй вопросы для самоанализа, рефлексии или уточняющие вопросы
- Цель — напомнить, что делать сегодня, и замотивировать

Стиль: живой, тёплый, не формальный. Не больше 7 строк суммарно."""
    return await call_claude([{"role": "user", "content": prompt}], system,
                             model=MODEL_SMART, max_tokens=MAX_TOKENS_NOTIF)

async def _build_evening_text(uid: int, profile: dict, local_now: datetime) -> str:
    address = profile.get("address") or profile.get("name") or ""
    tasks = get_tasks(uid)
    today_str = local_now.date().isoformat()
    done_today = db_fetch("""SELECT text FROM tasks WHERE user_id=? AND done=1
                              AND done_at >= ?""", (uid, today_str))
    done_block = "\n".join([f"• {t[0]}" for t in done_today[:5]]) if done_today else "Нет данных"
    system = build_system(profile, uid=uid)
    prompt = f"""Сгенерируй вечернее уведомление для пользователя.

Обращение: {address}
Выполнено сегодня: {done_block}
Открытых задач осталось: {len(tasks)}

Структура (одним цельным сообщением):
1. Тёплое вечернее приветствие
2. Короткий итог дня — что сделано (обобщённо, поддерживающе, не списком)
3. Одно намерение или фокус на завтра
4. «Программирование на удачу» — короткая вдохновляющая фраза-установка, которая настраивает на успех завтра (спокойно, по-взрослому, без пафоса)

ВАЖНО:
- НЕ задавай никаких вопросов пользователю
- Не спрашивай «что дал день» и не зови к рефлексии — пользователю этого не нужно в приветствии
- Не упоминай «сферы без внимания»

Стиль: мягкий, заботливый, человечный. Не более 6 строк суммарно."""
    return await call_claude([{"role": "user", "content": prompt}], system,
                             model=MODEL_SMART, max_tokens=MAX_TOKENS_NOTIF)

async def morning(context):
    utc_now = datetime.now(timezone.utc)
    users = db_fetch("SELECT user_id, profile FROM users WHERE onboarding_done=1")
    # Рассылка с пэйсингом: Telegram бан-лимит ~30 сообщений/сек в сумме.
    # 0.3с между юзерами = ~3 rps — безопасный режим даже на тысячах аккаунтов.
    PACE_SEC = 0.3
    sent_count = 0
    skipped_count = 0
    for uid, pj in users:
        profile = json.loads(pj)
        local_now = utc_now + timedelta(hours=get_user_tz_offset(profile))
        if local_now.hour != 8:
            continue
        today_local = local_now.date().isoformat()
        if _already_sent_today(profile, "last_morning_sent", today_local):
            skipped_count += 1
            continue
        try:
            response = await _build_morning_text(uid, profile, local_now)
            await context.bot.send_message(uid, response, parse_mode="Markdown")
            save_msg(uid, "assistant", response)
            _mark_sent(uid, "last_morning_sent", today_local)
            sent_count += 1
            logging.info(f"Morning sent to {uid} (local_date={today_local})")
        except Exception as e:
            logging.error(f"Morning notif error {uid}: {e}")
        await asyncio.sleep(PACE_SEC)
    if sent_count or skipped_count:
        logging.info(f"Morning summary: sent={sent_count}, skipped_dup={skipped_count}")

async def evening(context):
    utc_now = datetime.now(timezone.utc)
    users = db_fetch("SELECT user_id, profile FROM users WHERE onboarding_done=1")
    PACE_SEC = 0.3
    sent_count = 0
    skipped_count = 0
    for uid, pj in users:
        profile = json.loads(pj)
        local_now = utc_now + timedelta(hours=get_user_tz_offset(profile))
        if local_now.hour != 19:
            continue
        today_local = local_now.date().isoformat()
        if _already_sent_today(profile, "last_evening_sent", today_local):
            skipped_count += 1
            continue
        try:
            response = await _build_evening_text(uid, profile, local_now)
            await context.bot.send_message(uid, response, parse_mode="Markdown")
            save_msg(uid, "assistant", response)
            _mark_sent(uid, "last_evening_sent", today_local)
            sent_count += 1
            logging.info(f"Evening sent to {uid} (local_date={today_local})")
        except Exception as e:
            logging.error(f"Evening notif error {uid}: {e}")
        await asyncio.sleep(PACE_SEC)
    if sent_count or skipped_count:
        logging.info(f"Evening summary: sent={sent_count}, skipped_dup={skipped_count}")

async def weekly_review(context):
    utc_now = datetime.now(timezone.utc)
    users = db_fetch("SELECT user_id, profile FROM users WHERE onboarding_done=1")
    PACE_SEC = 0.5
    for uid, pj in users:
        profile = json.loads(pj)
        local_now = utc_now + timedelta(hours=get_user_tz_offset(profile))
        if local_now.hour != 10 or local_now.weekday() != 6:
            continue
        today_local = local_now.date().isoformat()
        if _already_sent_today(profile, "last_weekly_sent", today_local):
            continue
        _mark_sent(uid, "last_weekly_sent", today_local)
        address = profile.get("address") or profile.get("name") or ""
        tasks = get_tasks(uid)
        goals = get_goals(uid)
        frozen = get_frozen_items(uid)
        stats = get_sphere_stats(uid)

        goals_block = "\n".join([f"• {g[1]} — {g[4]}%" for g in goals[:6]]) if goals else "Целей нет"
        stats_block = "\n".join([f"• {SPHERES.get(k,k)}: {v} дн." for k,v in stats.items()]) if stats else "Нет данных"
        frozen_block = "\n".join([f"• {f[1]}" for f in frozen[:3]]) if frozen else ""

        system = build_system(profile, uid=uid)
        prompt = f"""Сгенерируй еженедельный обзор для пользователя.

Обращение: {address}
Неделя: {(local_now - timedelta(days=6)).strftime('%d.%m')} — {local_now.strftime('%d.%m.%Y')}
Открытых задач: {len(tasks)}
Прогресс целей:
{goals_block}
Активность по сферам за неделю:
{stats_block}
{f"Давно без движения:{chr(10)}{frozen_block}" if frozen_block else ""}

Структура:
1. Приветствие с неделей
2. Итоги по сферам — что активно, что игнорируется
3. Прогресс целей — честно и поддерживающе
4. Планы и фокус на следующую неделю
5. Если есть замороженные — мягко поднять
6. Завершение — вдохновляющее

Стиль: глубже обычного, аналитично но по-человечески. Не более 12 строк."""

        try:
            response = await call_claude([{"role": "user", "content": prompt}], system,
                                         model=MODEL_SMART, max_tokens=MAX_TOKENS_REVIEW)

            # Отправляем график перед текстом
            chart = generate_sphere_chart(uid)
            if chart:
                await context.bot.send_photo(uid, photo=chart,
                    caption="📊 Прогресс по сферам и целям за неделю")

            await context.bot.send_message(uid, response, parse_mode="Markdown")
            save_msg(uid, "assistant", response)

            # Предлагаем PDF
            await context.bot.send_message(uid, "Хочешь подробный PDF отчёт? Напиши /report")
        except Exception as e:
            logging.error(f"Weekly review error {uid}: {e}")
        await asyncio.sleep(PACE_SEC)

async def check_followup(context):
    pending = get_pending_followups()
    for uid, asked_at, attempts in pending:
        profile = get_profile(uid)
        history = get_history(uid, limit=6)
        system = build_system(profile, uid=uid)
        try:
            response = await call_claude(
                history + [{"role": "user", "content":
                    "Я не ответил на твой последний вопрос. Переформулируй его иначе — коротко, с другой стороны. "
                    "Не упоминай что я молчал."}],
                system, model=MODEL_SMART, max_tokens=MAX_TOKENS_DEFAULT)
            clean = await process_response(uid, response)
            await context.bot.send_message(uid, clean, parse_mode="Markdown")
            db_exec("UPDATE followup_queue SET asked_at=?, attempts=? WHERE user_id=?",
                    (datetime.now().isoformat(), attempts + 1, uid))
        except Exception as e:
            logging.error(f"Followup error {uid}: {e}")
        await asyncio.sleep(0.3)

async def cmd_report(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    await update.message.reply_text("Генерирую отчёт... 📊")
    chart = generate_sphere_chart(uid)
    if chart:
        await context.bot.send_photo(uid, photo=chart, caption="📊 Прогресс по сферам и целям")
    else:
        await update.message.reply_text("Нет данных для графика — добавь задачи и цели.")
    pdf = generate_pdf_report(uid)
    if pdf:
        await context.bot.send_document(uid, document=pdf, filename="nova_report.pdf",
                                        caption="📄 Подробный отчёт")
    else:
        await update.message.reply_text("PDF недоступен (библиотека reportlab не установлена).",
                                        reply_markup=main_keyboard())

from aiohttp import web

async def oauth_callback(request):
    code = request.rel_url.query.get("code")
    state = request.rel_url.query.get("state")
    if not code or not state:
        return web.Response(text="Ошибка авторизации — нет кода или state")
    try:
        import asyncio
        uid = int(state)
        if not GOOGLE_CLIENT_ID or not GOOGLE_CLIENT_SECRET:
            await request.app["bot"].send_message(uid,
                "❌ Calendar не настроен: отсутствуют GOOGLE_CLIENT_ID или GOOGLE_CLIENT_SECRET в переменных Railway.\n\nНужно добавить их в Railway → Variables.")
            return web.Response(text="Ошибка: Google credentials не настроены на сервере.")
        flow = get_oauth_flow()
        await asyncio.get_running_loop().run_in_executor(
            None, lambda: flow.fetch_token(code=code)
        )
        creds = flow.credentials
        save_google_token(uid, creds)

        # Показываем что уже есть в календаре
        try:
            events = await list_calendar_events(uid, max_results=25, include_past=False)
            cals = await list_calendars(uid)
            lines = ["✅ Google Календарь подключён!\n"]
            if len(cals) > 1:
                lines.append(f"📁 Найдено {len(cals)} календарей:")
                for c in cals[:5]:
                    marker = " _(основной)_" if c.get("primary") else ""
                    lines.append(f"  • {c['summary']}{marker}")
                lines.append("")
            if events:
                lines.append(f"📅 Предстоящих событий: *{len(events)}*")
                for e in events[:8]:
                    start = e['start'][:10] if e['start'] else "?"
                    lines.append(f"  • {e['summary']} — {start}")
                if len(events) > 8:
                    lines.append(f"  _...и ещё {len(events)-8}_")
                lines.append("\nХочешь изучить или удалить лишнее? Просто напиши мне.")
            else:
                lines.append("Событий пока нет — жду задачи для планирования 📅")
            await request.app["bot"].send_message(uid, "\n".join(lines), parse_mode="Markdown")
        except Exception as e:
            logging.error(f"Calendar show after oauth: {e}")
            await request.app["bot"].send_message(uid,
                "✅ Google Календарь успешно подключён!\n\nТеперь задачи будут появляться в календаре при планировании 📆")

        return web.Response(text="✅ Готово! Можешь закрыть эту вкладку и вернуться в Telegram.")
    except Exception as e:
        logging.error(f"OAuth callback error: {e}")
        try:
            uid = int(state)
            await request.app["bot"].send_message(uid,
                f"❌ Ошибка подключения Calendar:\n`{str(e)[:300]}`\n\nОтправь этот текст разработчику.")
        except Exception:
            pass
        return web.Response(text=f"Ошибка: {e}")

def main():
    init_db()
    app = Application.builder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start",   cmd_start))
    app.add_handler(CommandHandler("done",    cmd_done))
    app.add_handler(CommandHandler("help",    cmd_help))
    app.add_handler(CommandHandler("today",   cmd_today))
    app.add_handler(CommandHandler("week",    cmd_week))
    app.add_handler(CommandHandler("month",   cmd_month))
    app.add_handler(CommandHandler("goals",   cmd_goals))
    app.add_handler(CommandHandler("ideas",   cmd_ideas))
    app.add_handler(CommandHandler("focus",   cmd_focus))
    app.add_handler(CommandHandler("checkin", cmd_checkin))
    app.add_handler(CommandHandler("sphere",  cmd_sphere))
    app.add_handler(CommandHandler("reflect", cmd_reflect))
    app.add_handler(CommandHandler("wins",    cmd_wins))
    app.add_handler(CommandHandler("mood",    cmd_mood))
    app.add_handler(CommandHandler("habits",  cmd_habits))
    app.add_handler(CommandHandler("energy",  cmd_energy))
    app.add_handler(CommandHandler("journal", cmd_journal))
    app.add_handler(CommandHandler("brain",   cmd_brain))
    app.add_handler(CommandHandler("ask",     cmd_ask))
    app.add_handler(CommandHandler("review",  cmd_review))
    app.add_handler(CommandHandler("calendar",cmd_calendar))
    app.add_handler(CommandHandler("calreset",cmd_calreset))
    app.add_handler(CommandHandler("calshow", cmd_calshow))
    app.add_handler(CommandHandler("calinfo",cmd_calinfo))
    app.add_handler(CommandHandler("report",  cmd_report))
    app.add_handler(CommandHandler("settings",cmd_settings))
    app.add_handler(CommandHandler("profile", cmd_profile))
    app.add_handler(CommandHandler("plan",    cmd_plan))
    # Новые команды: финансы, GDPR, админ, парковка (доступны всегда)
    app.add_handler(CommandHandler("finance",       cmd_finance))
    app.add_handler(CommandHandler("export",        cmd_export))
    app.add_handler(CommandHandler("delete_me",     cmd_delete_me))
    app.add_handler(CommandHandler("myid",          cmd_myid))
    app.add_handler(CommandHandler("admin",         cmd_admin))
    app.add_handler(CommandHandler("parking",       cmd_parking))
    app.add_handler(CommandHandler("city",          cmd_city))
    app.add_handler(CommandHandler("forget",        cmd_forget))
    app.add_handler(CommandHandler("invite",        cmd_invite))
    app.add_handler(CommandHandler("applycode",     cmd_applycode))
    app.add_handler(CommandHandler("createpromo",   cmd_createpromo))
    app.add_handler(CommandHandler("draw",          cmd_draw))
    app.add_handler(CommandHandler("presentation",  cmd_presentation))
    app.add_handler(CommandHandler("backup",        cmd_backup))
    app.add_handler(CommandHandler("test_morning",  cmd_test_morning))
    app.add_handler(CommandHandler("test_evening",  cmd_test_evening))
    # Команды и хендлеры монетизации — регистрируем только если PAYMENTS_ENABLED.
    # Код cmd_subscribe и Stars-handlers сохранён в файле на будущее.
    if PAYMENTS_ENABLED:
        app.add_handler(CommandHandler("subscribe", cmd_subscribe))
        from telegram.ext import PreCheckoutQueryHandler
        app.add_handler(PreCheckoutQueryHandler(handle_pre_checkout))
        app.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, handle_successful_payment))
        logging.info("Payments: ENABLED")
    else:
        logging.info("Payments: disabled (set PAYMENTS_ENABLED=true to enable)")
    # Скрытые команды — не показываются в меню BotFather
    app.add_handler(CommandHandler("reset",   cmd_reset))
    app.add_handler(CommandHandler("newuser", cmd_newuser))
    app.add_handler(CallbackQueryHandler(handle_callback))
    app.add_handler(MessageHandler(filters.VOICE, handle_voice))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))
    jq = app.job_queue
    jq.run_repeating(morning, interval=3600, first=60)
    jq.run_repeating(evening, interval=3600, first=120)
    jq.run_repeating(weekly_review, interval=3600, first=180)
    jq.run_repeating(check_followup, interval=1800, first=300)
    # NPS-опрос: проверяем раз в сутки
    jq.run_repeating(check_feedback, interval=86400, first=600)

    from telegram import BotCommand

    BOT_COMMANDS = [
        BotCommand("start",    "Знакомство и онбординг"),
        BotCommand("help",     "Что умею и как работать"),
        BotCommand("today",    "Задачи на сегодня"),
        BotCommand("week",     "План на неделю"),
        BotCommand("month",    "План на месяц"),
        BotCommand("goals",    "Цели и прогресс"),
        BotCommand("ideas",    "Идеи и желания"),
        BotCommand("focus",    "Главная задача прямо сейчас"),
        BotCommand("checkin",  "Чекин состояния и энергии"),
        BotCommand("sphere",   "Колесо жизни"),
        BotCommand("reflect",  "Вопрос для самоанализа"),
        BotCommand("wins",     "Победы и достижения"),
        BotCommand("mood",     "Трекер настроения"),
        BotCommand("habits",   "Трекер привычек"),
        BotCommand("energy",   "Отметить уровень энергии"),
        BotCommand("journal",  "Личный дневник"),
        BotCommand("brain",    "Выгрузка мыслей"),
        BotCommand("ask",      "Честный анализ от Новы"),
        BotCommand("review",   "Ежемесячный разбор"),
        BotCommand("calendar", "Подключить Google Календарь"),
        BotCommand("calshow",  "Показать события в Calendar"),
        BotCommand("report",   "Отчёт, графики, PDF"),
        BotCommand("settings", "Настройки"),
        BotCommand("profile",  "Мой профиль"),
        BotCommand("plan",     "Недельный план — задачи + приоритеты"),
        BotCommand("finance",      "Мои траты (распознаю по фото чека)"),
        BotCommand("parking",      "Вернуться к отложенным темам из знакомства"),
        BotCommand("city",         "Указать свой город"),
        BotCommand("invite",       "Моя реферальная ссылка"),
        BotCommand("applycode",    "Применить промокод"),
        BotCommand("draw",         "Сгенерировать картинку"),
        BotCommand("presentation", "Сделать презентацию .pptx"),
        BotCommand("forget",       "Забыть разговор или память"),
        BotCommand("export",       "Скачать все свои данные"),
        BotCommand("delete_me",    "Удалить все мои данные"),
    ]
    if PAYMENTS_ENABLED:
        BOT_COMMANDS.append(BotCommand("subscribe", "Мой тариф и подписка"))

    async def start_web(app_obj):
        await app_obj.bot.set_my_commands(BOT_COMMANDS)
        logging.info("Bot commands registered")
        web_app = web.Application()
        web_app["bot"] = app_obj.bot
        web_app.router.add_get("/oauth/callback", oauth_callback)
        runner = web.AppRunner(web_app)
        await runner.setup()
        site = web.TCPSite(runner, "0.0.0.0", int(os.environ.get("PORT", 8080)))
        await site.start()
        logging.info("Web server started")

    app.post_init = start_web
    logging.info("Nova is running...")
    app.run_polling()

if __name__ == "__main__":
    main()